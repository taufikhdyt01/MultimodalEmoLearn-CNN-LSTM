"""
Add slide "Pertanyaan Bimbingan — Scope Eksperimen Tesis" untuk konsultasi
ke Pak Fitra: list eksperimen yang sudah dilakukan + yang partial + yang
tidak memungkinkan dalam waktu tersisa.

Posisi: setelah slide 141 (Pertanyaan Paper) dan sebelum slide 142
(Update Komunikasi Pak Fitra). Push Terimakasih ke index akhir.
"""
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

PPTX = Path("docs/PPT Bimbingan.pptx")
BACKUP = Path("docs/PPT Bimbingan_pre_scope_tesis.pptx")


def add_textbox(slide, left, top, width, height, lines, *, color=None,
                fill_color=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    if fill_color:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill_color
    for i, item in enumerate(lines):
        if len(item) == 3:
            txt, size_pt, bold = item
            line_color = color
        else:
            txt, size_pt, bold, line_color = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i == 0:
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
            p.text = ""
        run = p.add_run()
        run.text = txt
        if size_pt is not None:
            run.font.size = Pt(size_pt)
        if bold is not None:
            run.font.bold = bold
        if line_color:
            run.font.color.rgb = line_color
    return tb


# Color palette
GREEN = RGBColor(0x1B, 0x7A, 0x3E)
ORANGE = RGBColor(0xC4, 0x6B, 0x10)
RED = RGBColor(0xB0, 0x2A, 0x2A)
GRAY = RGBColor(0x55, 0x55, 0x55)
BLUE = RGBColor(0x2E, 0x5E, 0x96)


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    shutil.copy2(PPTX, BACKUP)
    print(f"Backup: {BACKUP}")

    prs = Presentation(PPTX)
    n_before = len(prs.slides)
    print(f"Before: {n_before} slides")

    # Locate target positions:
    # idx 140 (slide 141): Pertanyaan Bimbingan — Update Paper
    # idx 141 (slide 142): Update Komunikasi ke Pak Fitra
    # idx 142 (slide 143): Terimakasih
    # Insert new slide at idx 141 (becomes slide 142), pushing others +1.
    insert_at = 141  # 0-based index
    # Use slide 141 (Pertanyaan Paper) layout as template
    src = prs.slides[140]
    new_slide = prs.slides.add_slide(src.slide_layout)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    new_sld = slides_list[-1]
    xml_slides.remove(new_sld)
    xml_slides.insert(insert_at, new_sld)

    # ── Title ──
    add_textbox(new_slide, 0.30, 0.08, 9.40, 0.40,
                [("Pertanyaan Bimbingan — Scope Eksperimen Tesis", 18, True)])

    # ── Subtitle / pertanyaan ──
    add_textbox(new_slide, 0.30, 0.52, 9.40, 0.32,
                [("Mohon arahan: dari list eksperimen di bawah, mana yang sebaiknya dibahas "
                  "di BAB Hasil tesis, dan apakah yang masih partial/pending masih perlu dikejar?",
                  10, False, BLUE)])

    # ─── KOLOM 1: DONE (akan dibahas) ───
    col1_x = 0.30
    col1_w = 4.80
    add_textbox(new_slide, col1_x, 0.92, col1_w, 0.30,
                [("✅ DONE — Layak Dibahas (15)", 12, True, GREEN)])

    done_lines = [
        ("Pipeline + Dataset Primer", 9, True, None),
        ("• Preprocessing 224×224, MediaPipe 68-pt landmark, Face API conf≥60%", 8, False, None),
        ("• Split per-subject 78/8.5/13.5%, total 6,795 sampel front-only", 8, False, None),
        ("Grid 7-class (27 configs, BAB 4)", 9, True, None),
        ("• 9 arch × 3 scenario, juara: Early Fusion TL B3 macro=0.333", 8, False, None),
        ("Grid 3-class reframe (27 configs, BAB 4) ⭐", 9, True, None),
        ("• Russell 1980 valence — juara overall: Late Fusion TL B3 val=0.623", 8, False, None),
        ("4-class (archived, narasi reframe)", 9, True, None),
        ("• Dijadikan motivasi reframe ke 3-class (literature gap)", 8, False, None),
        ("Soft Label Training (negative)", 9, True, None),
        ("• nb 71/72/78 — KL-div, BCE, focal — no improvement", 8, False, None),
        ("GradCAM Observasi (nb 73)", 9, True, None),
        ("Space Alignment / CCA (nb 75) — positive", 9, True, None),
        ("• Top-5 alignment = 0.978, image+landmark co-encoded", 8, False, None),
        ("Confidence-Stratified Analysis", 9, True, None),
        ("• Refutes label-noise hypothesis (slide 39)", 8, False, None),
        ("CBAM Attention (nb 80) — negative", 9, True, None),
        ("• 4 configs, none beat LF TL B3 val=0.623", 8, False, None),
        ("Geometric Features Liliana 2019 (nb 81)", 9, True, None),
        ("• 20-d GF — negative, tesis-only (3-class)", 8, False, None),
        ("Expert Validation Pak Dephilia", 9, True, None),
        ("• 128 entries divalidasi, Cohen's κ=0.45 (moderate)", 8, False, None),
        ("Cross-dataset benchmark 7-class", 9, True, None),
        ("• RAF-DB, KDEF, cross-dataset to/from Primer", 8, False, None),
    ]
    add_textbox(new_slide, col1_x, 1.22, col1_w, 4.10, done_lines)

    # ─── KOLOM 2: PARTIAL / OUT OF SCOPE ───
    col2_x = 5.20
    col2_w = 4.50

    # Partial / Pending
    add_textbox(new_slide, col2_x, 0.92, col2_w, 0.30,
                [("🔄 PARTIAL / Pending VPS (3)", 12, True, ORANGE)])
    partial_lines = [
        ("3-class Skema 1 benchmark (nb 84)", 9, True, None),
        ("• CK+ ✅ + JAFFE ✅ done", 8, False, None),
        ("• RAF-DB, KDEF — pending VPS recovery", 8, False, ORANGE),
        ("3-class Skema 2 cross-dataset (nb 85)", 9, True, None),
        ("• Train Primer 3c → test RAF-DB/KDEF/CK+/JAFFE — pending VPS", 8, False, ORANGE),
        ("History logging full grid (nb 83)", 9, True, None),
        ("• Hanya 5 config dilakukan — full 27 belum logged", 8, False, ORANGE),
    ]
    add_textbox(new_slide, col2_x, 1.22, col2_w, 1.80, partial_lines)

    # Not feasible
    add_textbox(new_slide, col2_x, 3.05, col2_w, 0.30,
                [("❌ TIDAK MEMUNGKINKAN — Out of Scope (5)", 12, True, RED)])
    notfeasible_lines = [
        ("Pitaloka 2017 GCN Ablation", 9, True, None),
        ("• Tidak dilakukan — alasan: implementasi GCN dari awal "
         "(graph construction dari landmark) butuh ≥1 minggu", 8, False, None),
        ("Liliana 2019 FEIS Fuzzy Rule", 9, True, None),
        ("• Non-DL baseline — manual rule design + tuning, "
         "estimasi 2-3 hari + di luar fokus DL", 8, False, None),
        ("Ghost / Triplet loss approach", 9, True, None),
        ("• Tidak dilakukan — kompleksitas pair sampling + class weight interaksi", 8, False, None),
        ("CBAM / Geometric untuk 7-class scheme", 9, True, None),
        ("• Hanya diuji di 3-class (juara scheme); pattern negative di 3c → "
         "kemungkinan lebih buruk di kelas minoritas 7c", 8, False, None),
        ("Multitask learning (valence + arousal)", 9, True, None),
        ("• Tidak ada label arousal di Primer dataset (Face API hanya 7 emotion)", 8, False, None),
    ]
    add_textbox(new_slide, col2_x, 3.35, col2_w, 2.00, notfeasible_lines)

    # ── Footer pertanyaan ──
    add_textbox(new_slide, 0.30, 5.30, 9.40, 0.32,
                [("Pertanyaan: (1) item DONE mana yang prioritas highlight di tesis? "
                  "(2) RAF-DB/KDEF 3c benchmark perlu dikejar atau cukup CK+/JAFFE? "
                  "(3) Dari yang TIDAK MEMUNGKINKAN, ada yang Bapak prioritaskan untuk dipaksakan?",
                  9, True, BLUE)])

    prs.save(PPTX)
    print(f"After:  {len(prs.slides)} slides (+1 = 'Pertanyaan Scope Eksperimen Tesis')")
    print(f"Saved:  {PPTX}")


if __name__ == "__main__":
    main()
