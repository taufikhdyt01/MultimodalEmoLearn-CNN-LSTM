"""
Insert slide baru "SLIDE 36: Motivasi Reframe 3-Class" sebelum SLIDE 36 Results.

Content: arahan Pak Fitra + hasil literature search (Savchenko 2022 VGAF) yang
justify pindah dari 4-class ke 3-class polarity (Russell circumplex).

Position:
  132 (1-idx): SLIDE 35 Space Alignment
  133 (new):   SLIDE 36: Motivasi Reframe 3-Class  ← insert
  134 (was 133): SLIDE 36 (lanjutan): 3-Class Results
  135 (was 134): SLIDE 36 (lanjutan 2): 3-Class Best Model CM
  136 (was 135): Rencana
  137 (was 136): Diskusi
  138 (was 137): Pertanyaan Bimbingan Paper
  139 (was 138): Terimakasih
"""
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

PPTX   = Path('docs/PPT Bimbingan.pptx')
BACKUP = Path('docs/PPT Bimbingan_pre_motivasi.pptx')


def add_textbox(slide, left, top, width, height, lines, color=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, (txt, size_pt, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i == 0:
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
            p.text = ''
        run = p.add_run()
        run.text = txt
        if size_pt is not None: run.font.size = Pt(size_pt)
        if bold is not None: run.font.bold = bold
        if color: run.font.color.rgb = color
    return tb


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    shutil.copy2(PPTX, BACKUP)
    print(f'Backup: {BACKUP}')

    prs = Presentation(PPTX)
    # SLIDE 35 Space Alignment is at 0-idx 131 (1-idx 132)
    # Insert new motivasi slide at 0-idx 132 (becomes 1-idx 133)
    # Existing 3-class results slide shifts from 0-idx 132 → 0-idx 133 (1-idx 134)

    src = prs.slides[131]  # layout source
    new_slide = prs.slides.add_slide(src.slide_layout)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    xml_slides.remove(slides_list[-1])
    xml_slides.insert(132, slides_list[-1])

    # ── Content ──

    # Title
    add_textbox(new_slide, 0.3, 0.08, 9.4, 0.40, [
        ('SLIDE 36: Reframe ke 3-Class — Motivasi & Literature Justification', 17, True),
    ])

    # Section 1: Arahan Pak Fitra
    add_textbox(new_slide, 0.3, 0.55, 9.4, 0.30, [
        ('1. Arahan Pak Fitra', 12, True),
    ], color=RGBColor(0x2E, 0x5E, 0x96))
    add_textbox(new_slide, 0.3, 0.88, 9.4, 0.35, [
        ('"Cari precedent literatur untuk skema 4-class yang dipakai (neutral/happy/sad/negative)"',
         10, False),
    ])

    # Section 2: Hasil Literature Search
    add_textbox(new_slide, 0.3, 1.35, 9.4, 0.30, [
        ('2. Hasil Literature Search', 12, True),
    ], color=RGBColor(0x2E, 0x5E, 0x96))
    add_textbox(new_slide, 0.3, 1.68, 9.4, 0.85, [
        ('• Skema 4-class (neutral/happy/sad/negative) — sulit ditemukan precedent di literatur FER '
         'konteks online learning. Mapping ini arbitrary, tidak match standard taxonomy.', 9.5, False),
        ('• Skema 3-class polarity (positive/neutral/negative) — banyak digunakan, '
         'sesuai Russell 1980 circumplex valence dimension. Standard di affective computing.', 9.5, False),
    ])

    # Section 3: Key Paper Citation
    add_textbox(new_slide, 0.3, 2.68, 9.4, 0.30, [
        ('3. Paper Paling Relevan — Savchenko et al. (2022)', 12, True),
    ], color=RGBColor(0x2E, 0x5E, 0x96))
    add_textbox(new_slide, 0.3, 3.00, 9.4, 1.40, [
        ('"Classifying Emotions and Engagement in Online Learning Based on a Single Facial '
         'Expression Recognition Neural Network"', 10, True),
        ('IEEE Transactions on Affective Computing, Vol. 13, No. 4 (2022)', 9, False),
        ('', 3, False),
        ('• Topik: FER untuk e-learning engagement — persis sama dengan penelitian ini', 9.5, False),
        ('• Dataset: VGAF (Video Group Affect)', 9.5, False),
        ('• Skema: 3-class polarity (positive / neutral / negative)', 9.5, True),
        ('• Justifikasi empiris: menunjukkan 3-class polarity adalah konvensi untuk online learning FER', 9.5, False),
    ])

    # Section 4: Keputusan & Mapping
    add_textbox(new_slide, 0.3, 4.55, 9.4, 0.30, [
        ('4. Keputusan: Reframe Paper JITeCS ke 3-Class', 12, True),
    ], color=RGBColor(0x2E, 0x5E, 0x96))
    add_textbox(new_slide, 0.3, 4.88, 9.4, 0.60, [
        ('Mapping:  happy + surprised → positive   |   neutral → neutral   |   '
         'sad + angry + fearful + disgusted → negative', 9.5, True),
        ('Training ulang dengan label 3-class sudah dijalankan → nb 79 (next slide). '
         'Hasil: imbalance 1:14 (vs 7c 1:1138), best Macro F1 0.623 (val-tuned).',
         9, False),
    ])

    prs.save(PPTX)
    print(f'After: {len(prs.slides)} slides (inserted motivasi at 1-idx 133)')


if __name__ == '__main__':
    main()
