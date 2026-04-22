"""
Update PPT Bimbingan:
  (A) Insert SLIDE 33 Soft Label Training (nb 71) sebagai slide baru setelah slide 126.
  (B) Rewrite slide 127 (Rencana Eksperimen Lanjutan) dengan plan 8 eksplorasi lanjutan
      (4 arahan dosen + 4 turunan paper ref). Soft Label ditandai DONE.
  (C) Update slide 128 (Diskusi & Konsultasi) dengan angka val-tuned + pertanyaan terkini.

Input:  docs/PPT Bimbingan.pptx
Output: docs/PPT Bimbingan.pptx (in-place, backup ke _pre_eksplorasi.pptx)

After run, slide order menjadi:
  126: (lama) Skema 2 RAF-DB/KDEF
  127: (BARU) SLIDE 33 Soft Label Training
  128: (rewrite, was 127) Rencana Eksplorasi Lanjutan
  129: (update, was 128) Diskusi & Konsultasi
  130: (was 129) Terimakasih
"""
import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

PPTX = Path('docs/PPT Bimbingan.pptx')
BACKUP = Path('docs/PPT Bimbingan_pre_eksplorasi.pptx')

# ───────── Helpers ─────────

def set_text_preserve(text_frame, lines, style_from_run=None):
    """Replace text_frame content with lines, trying to preserve style from first run."""
    # Save reference style
    ref_font = None
    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
        ref_font = text_frame.paragraphs[0].runs[0].font
    # Clear all existing paragraphs except first
    for para in list(text_frame.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    # Clear runs in first paragraph
    first_para = text_frame.paragraphs[0]
    for run in list(first_para.runs):
        run._r.getparent().remove(run._r)
    first_para.text = ''
    # Add lines
    for i, (txt, size_pt, bold) in enumerate(lines):
        if i == 0:
            p = first_para
        else:
            p = text_frame.add_paragraph()
        run = p.add_run()
        run.text = txt
        if size_pt is not None:
            run.font.size = Pt(size_pt)
        if bold is not None:
            run.font.bold = bold


def set_cell_text_preserve(cell, new_text):
    tf = cell.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        first_run = tf.paragraphs[0].runs[0]
        first_run.text = new_text
        for run in list(tf.paragraphs[0].runs)[1:]:
            run._r.getparent().remove(run._r)
        for para in list(tf.paragraphs)[1:]:
            para._p.getparent().remove(para._p)
    else:
        cell.text = new_text


def duplicate_slide(prs, src_idx, target_pos):
    """Clone slide at src_idx, move to target_pos (0-indexed). Return new slide."""
    src = prs.slides[src_idx]
    new_slide = prs.slides.add_slide(src.slide_layout)
    # Remove placeholder shapes from new layout
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
    # Deep-copy each shape from source
    for shape in src.shapes:
        new_el = deepcopy(shape._element)
        new_slide.shapes._spTree.append(new_el)
    # Move new slide (currently last) to target_pos
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    xml_slides.remove(slides_list[-1])
    xml_slides.insert(target_pos, slides_list[-1])
    return new_slide


def replace_table_row(table, ridx, cell_texts):
    row = table.rows[ridx]
    for i, t in enumerate(cell_texts):
        if i < len(row.cells):
            set_cell_text_preserve(row.cells[i], t)


def clone_and_append_row(table, cell_texts):
    tbl = table._tbl
    new_tr = deepcopy(tbl.tr_lst[-1])
    tbl.append(new_tr)
    new_row = table.rows[len(table.rows) - 1]
    for i, t in enumerate(cell_texts):
        if i < len(new_row.cells):
            set_cell_text_preserve(new_row.cells[i], t)


def drop_extra_table_rows(table, keep_count):
    """Remove rows beyond keep_count."""
    tbl = table._tbl
    rows = tbl.tr_lst
    while len(rows) > keep_count:
        tbl.remove(rows[-1])
        rows = tbl.tr_lst


# ───────── (A) Soft Label slide content ─────────

def populate_soft_label_slide(slide):
    """slide was cloned from slide 120 (SLIDE 30 lanjutan):
       shapes: [0]=? [1]=title [2]=subtitle [3]=table 6r×4c [4]=Temuan hdr [5-8]=temuan boxes
    """
    shapes = list(slide.shapes)
    # Shape 1: title
    set_text_preserve(shapes[1].text_frame, [
        ('SLIDE 33: Soft Label Training (Eksplorasi Liliana 2019-inspired, nb 71)', 18, True)
    ])
    # Shape 2: subtitle
    set_text_preserve(shapes[2].text_frame, [
        ('4 Loss Variants — CNN TL 4-class B1 Baseline (isolate efek loss function)', 11, False)
    ])
    # Shape 3: 6-row × 4-col table. Convert to 5 rows × 6 cols? Keep 4 cols for simplicity:
    # Row 0 header: Config | Macro F1 | Micro F1 | Weighted F1 (we keep 4 cols)
    # Rows 1-4: Hard CE, Soft CE, KL-div, Label Smoothing
    # Row 5: reference note
    tbl = shapes[3].table
    # Shrink to 6 rows (should already be 6)
    drop_extra_table_rows(tbl, 6)
    replace_table_row(tbl, 0, ['Config', 'Macro F1', 'Micro F1', 'W-F1'])
    replace_table_row(tbl, 1, ['A. Hard CE (baseline)', '0.4269', '0.6986', '0.7155'])
    replace_table_row(tbl, 2, ['B. Soft CE', '0.4667', '0.8407', '0.8281'])
    replace_table_row(tbl, 3, ['C. KL-divergence  ⭐', '0.5170', '0.8213', '0.8258'])
    replace_table_row(tbl, 4, ['D. Label Smoothing (ε=0.1)', '0.4418', '0.8202', '0.8089'])
    replace_table_row(tbl, 5, ['Ref: CNN TL 4c B1 hard (nb 54) = 0.456', '', '', ''])

    # Shape 4: "Temuan" header
    set_text_preserve(shapes[4].text_frame, [('Temuan (41-44)', 12, True)])
    # Shapes 5-8: 4 Temuan boxes
    temuan = [
        ('Temuan 41: KL-divergence soft label beat Hard CE +0.090 Macro F1',
         'CNN TL 4c: 0.427 (hard) → 0.517 (KL-div). Signifikan. Konsisten hipotesis Liliana 2019: natural emotion recognition diuntungkan target fuzzy/distribusional.'),
        ('Temuan 42: Hierarki efektivitas — KL > Soft CE > Label Smoothing > Hard CE',
         'KL (Face API real) > Soft CE (sama distribusi, beda loss) > Label Smoothing (artificial ε=0.1) > Hard CE. Menunjukkan distribusi Face API membawa sinyal real, bukan sekadar regularisasi.'),
        ('Temuan 43: KL-div converge epoch 1 — perlu verifikasi',
         'best_epoch = 1 (stale 15 epoch lalu early stop). Perlu cek training history + multi-seed run untuk pastikan bukan fluke.'),
        ('Temuan 44: Soft label bump accuracy 0.70 → 0.84',
         'Hard CE acc 0.70 vs Soft CE acc 0.84. Model tidak selalu prediksi neutral dengan soft target — lebih "aware" kelas lain. Next: nb 72 extend ke Late Fusion TL 4c B3.'),
    ]
    for i, (heading, body) in enumerate(temuan):
        shp = shapes[5 + i]
        set_text_preserve(shp.text_frame, [
            (heading, 9.5, True),
            (body, 9, False),
        ])


# ───────── (B) Rewrite slide "Rencana Eksplorasi Lanjutan" ─────────

def rewrite_rencana_slide(slide):
    """Reuse 6 shapes of old Rencana slide.
       [0] title, [1] motivasi, [2] kol kiri, [3] kol kanan, [4] prioritas/wide, [5] status banner
    """
    shapes = list(slide.shapes)
    # [0] title
    set_text_preserve(shapes[0].text_frame, [
        ('Rencana & Status Eksplorasi Lanjutan (8 Arah)', 18, True)
    ])
    # [1] motivasi
    set_text_preserve(shapes[1].text_frame, [
        ('Best saat ini: Intermediate TL 4-class B3 = Macro F1 0.521 (val-tuned proper).', 10, True),
        ('8 arah eksplorasi disiapkan untuk memperkuat BAB eksperimen & novelty tesis — 4 arahan dosen + 4 turunan paper ref.', 9.5, False),
    ])
    # [2] kiri: Arahan Dosen (4)
    set_text_preserve(shapes[2].text_frame, [
        ('Arahan Dosen (4 item)', 11, True),
        ('', 4, False),
        ('★★★  1. GradCAM Evaluation — visualisasi fokus model di best Late Fusion TL + CNN TL (1-2 hari)', 9, False),
        ('★★    2. Space Alignment — CCA/t-SNE antara CNN (256-d) vs FCNN (128-d) features (1 hari)', 9, False),
        ('★★★★ 3. Attention Module — CBAM / Ghost / Triplet; target beat Late Fusion TL (4-6 hari)', 9, False),
        ('★★★★ 4. Geometric Features (Liliana 2019) — 20-dim GF dari 10 facial components (2-3 hari)', 9, False),
    ])
    # [3] kanan: Turunan Paper (4)
    set_text_preserve(shapes[3].text_frame, [
        ('Turunan Paper Referensi (4 item)', 11, True),
        ('', 4, False),
        ('★★★★★ 5. Soft Label Training ✅ DONE — KL-div 0.517 vs Hard CE 0.427 (+0.09), nb 71', 9, True),
        ('★★★★ 6. Fuzzy Rule-Based FEIS — replikasi Liliana 2019 non-DL baseline (2-3 hari)', 9, False),
        ('★★★★★ 7. Geometric + Soft Label Combined — double-extension Liliana (4-5 hari)', 9, False),
        ('★★    8. GCN Preprocessing Ablation — Pitaloka 2017 quick win (0.5-1 hari)', 9, False),
    ])
    # [4] prioritas wide
    set_text_preserve(shapes[4].text_frame, [
        ('Prioritas tertinggi (★★★★★): (5) Soft Label DONE, (7) Geometric + Soft Label Combined', 10, True),
        ('', 3, False),
        ('Next step konkret: notebook 72_soft_label_late_fusion.ipynb — extend soft label ke Late Fusion TL 4c + B3 (potensi beat 0.521).', 9.5, False),
        ('Detail lengkap + matriks eksperimen: docs/eksplorasi_lanjutan.md', 9, False),
    ])
    # [5] status
    set_text_preserve(shapes[5].text_frame, [
        ('Status: 1/8 DONE (Soft Label) | 7/8 planning — Total estimasi ~15-25 hari eksperimen jika semua diambil', 9.5, True)
    ])


# ───────── (C) Update Diskusi slide (tables) ─────────

def update_diskusi_slide(slide):
    """[0] title, [1] RQ table 4r×3c, [2] Konsultasi table 5r×3c"""
    shapes = list(slide.shapes)
    # Keep title
    # RQ table (4r × 3c): RQ | Pertanyaan | Jawaban
    rq_tbl = shapes[1].table
    replace_table_row(rq_tbl, 0, ['RQ', 'Pertanyaan Penelitian', 'Jawaban / Hasil'])
    replace_table_row(rq_tbl, 1, [
        'RQ1',
        'Apakah multimodal fusion outperforms single-modality?',
        'Ya. Best fusion (Intermediate TL 4c B3 = 0.521) > best single-modal (CNN TL 4c B3 = 0.507, FCNN 4c B2 = 0.459). Gain moderat ~+0.01-0.06 Macro F1.'
    ])
    replace_table_row(rq_tbl, 2, [
        'RQ2',
        'Strategi fusion mana yang terbaik (Early / Intermediate / Late)?',
        '4-class: Intermediate TL (0.521) > Late Fusion scratch B2 (0.479) ≈ Early Fusion TL (0.471). 7-class: Early Fusion TL B3 (0.333) > Intermediate TL (0.292). Val-tuned proper, no leakage.'
    ])
    replace_table_row(rq_tbl, 3, [
        'RQ3',
        'Transfer Learning effectiveness?',
        'TL konsisten unggul (Intermediate TL vs scratch: +0.127 B3 4c). Gain Late Fusion marginal post val-tuning karena w ≈ 0 pilih FCNN-only. Best 4c val-tuned: Intermediate TL > Late Fusion.'
    ])
    # Konsultasi table (5r × 3c)
    kt = shapes[2].table
    replace_table_row(kt, 0, ['No', 'Topik Konsultasi', 'Pertanyaan'])
    replace_table_row(kt, 1, [
        '1',
        'Best overall bergeser: Late Fusion TL (0.567 test-tuned) → Intermediate TL (0.521 val-tuned)',
        'Fix test-set leakage di Late Fusion grid-search w. Mohon konfirmasi narasi paper: pakai angka val-tuned proper (0.521) atau report both dengan caveat?'
    ])
    replace_table_row(kt, 2, [
        '2',
        'Soft Label Training (nb 71) — KL-divergence +0.09 Macro F1',
        'Hasil promising di CNN TL single (0.427 → 0.517). Lanjut ke nb 72 (Late Fusion + soft label + B3) untuk coba beat 0.521? Atau cukup dicantumkan sebagai eksplorasi?'
    ])
    replace_table_row(kt, 3, [
        '3',
        'Eksplorasi lanjutan: GradCAM, Space Alignment, Attention (CBAM/Ghost/Triplet), Geometric Features (Liliana)',
        'Prioritas mana yang diambil dulu untuk BAB eksperimen tesis? 4-6 item dalam ~2-3 minggu, atau fokus 1-2 item lebih dalam?'
    ])
    replace_table_row(kt, 4, [
        '4',
        'Scope paper JITeCS vs tesis',
        'Paper fokus Primer conf60 (5 fusion strategy + TL + imbalance). Eksplorasi lanjutan (8 item) masuk ke tesis saja, bukan paper?'
    ])


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    shutil.copy2(PPTX, BACKUP)
    print(f'Backup: {BACKUP}')

    prs = Presentation(PPTX)
    n_before = len(prs.slides)
    print(f'Before: {n_before} slides')

    # (A) Clone slide 120 (SLIDE 30 lanjutan template), insert at position 126 (0-idx)
    # → becomes new slide 127 (1-idx)
    print('\n(A) Cloning slide 120 (template) → insert at pos 126 (becomes slide 127)')
    new_soft = duplicate_slide(prs, src_idx=120 - 1, target_pos=126)
    populate_soft_label_slide(new_soft)
    print('    Soft Label slide populated')

    # After insert: slides shift by 1.
    # Old slide 127 (Rencana) is now at 1-idx 128 (0-idx 127)
    # Old slide 128 (Diskusi) is now at 1-idx 129 (0-idx 128)
    print('\n(B) Rewrite Rencana slide (now at 1-idx 128 / 0-idx 127)')
    rewrite_rencana_slide(prs.slides[127])

    print('\n(C) Update Diskusi slide (now at 1-idx 129 / 0-idx 128)')
    update_diskusi_slide(prs.slides[128])

    prs.save(PPTX)
    print(f'\nAfter: {len(prs.slides)} slides')
    print(f'Saved: {PPTX}')


if __name__ == '__main__':
    main()
