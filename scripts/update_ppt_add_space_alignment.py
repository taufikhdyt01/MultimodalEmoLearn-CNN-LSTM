"""
Add SLIDE 35 Space Alignment ke PPT sebelum Rencana Eksplorasi.

Slide baru di 1-idx 132 (setelah GradCAM slide 131), pushing:
  Rencana 132 → 133
  Diskusi 133 → 134
  Paper Q  134 → 135
  Terimakasih 135 → 136
"""
import shutil
from copy import deepcopy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

PPTX   = Path('docs/PPT Bimbingan.pptx')
BACKUP = Path('docs/PPT Bimbingan_pre_spacealign.pptx')

FIG_DIR = Path('docs/figures/space_alignment')
CCA_PNG = FIG_DIR / 'cca_correlations.png'
RET_PNG = FIG_DIR / 'retrieval_topk.png'


def add_textbox(slide, left, top, width, height, lines, color=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(lines):
        txt, size_pt, bold = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i == 0:
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
            p.text = ''
        run = p.add_run()
        run.text = txt
        if size_pt is not None:
            run.font.size = Pt(size_pt)
        if bold is not None:
            run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return tb


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    for p in [CCA_PNG, RET_PNG]:
        if not p.exists():
            raise FileNotFoundError(p)

    shutil.copy2(PPTX, BACKUP)
    print(f'Backup: {BACKUP}')

    prs = Presentation(PPTX)
    # GradCAM = 1-idx 131 = 0-idx 130; insert new after it (becomes 1-idx 132 = 0-idx 131)
    src = prs.slides[130]
    new_slide = prs.slides.add_slide(src.slide_layout)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    # Move to position 131 (0-idx), i.e. becomes 1-idx 132
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    xml_slides.remove(slides_list[-1])
    xml_slides.insert(131, slides_list[-1])

    # ── Populate content ──

    # Title
    add_textbox(new_slide, 0.3, 0.08, 9.4, 0.40, [
        ('SLIDE 35: Space Alignment (arahan dosen #2, nb 75)', 18, True),
    ])

    # Description
    add_textbox(new_slide, 0.3, 0.50, 9.4, 0.32, [
        ('Evaluasi alignment feature space CNN TL (256-d) vs FCNN (128-d) — 4 metode analitis '
         'di test set Primer n=929. Positive result, methodology-clean.', 9.5, False),
    ], color=RGBColor(0x55, 0x55, 0x55))

    # ── Figure 1: CCA ──
    new_slide.shapes.add_picture(str(CCA_PNG), Inches(0.25), Inches(0.90),
                                  width=Inches(4.60), height=Inches(2.25))
    add_textbox(new_slide, 0.25, 3.18, 4.60, 0.25, [
        ('CCA correlations — top-5 mean = 0.978', 9, True),
    ])

    # ── Figure 2: Retrieval ──
    new_slide.shapes.add_picture(str(RET_PNG), Inches(5.15), Inches(0.90),
                                  width=Inches(4.60), height=Inches(2.25))
    add_textbox(new_slide, 5.15, 3.18, 4.60, 0.25, [
        ('Cross-modal retrieval — top-5 = 95.3% (random = 0.54%)', 9, True),
    ])

    # ── Quantitative summary ──
    add_textbox(new_slide, 0.3, 3.52, 9.4, 0.75, [
        ('Metrik Kunci:', 9.5, True),
        ('• CCA top-5 mean = 0.978  (strong linear alignment, well above 0.5 threshold)', 9, False),
        ('• Paired cosine overall = 0.939  |  per-class konsisten 0.937-0.951 (termasuk minoritas negative n=8)', 9, False),
        ('• Cross-modal retrieval top-1 = 0.670, top-5 = 0.953, top-20 = 1.000  (random baseline top-5 = 0.005)', 9, False),
    ])

    # ── Interpretasi ──
    add_textbox(new_slide, 0.3, 4.32, 9.4, 0.65, [
        ('Interpretasi:', 9.5, True),
        ('Image dan landmark stream belajar representasi yang nyaris co-linear di latent space — justifikasi empiris '
         'untuk Intermediate Fusion (concat). Konsisten dengan Intermediate TL juara val-tuned 4c (0.521) > '
         'Late Fusion TL (0.466). Streams semantically paired meskipun modality berbeda.', 8.5, False),
    ])

    # ── Status ──
    add_textbox(new_slide, 0.3, 5.05, 9.4, 0.45, [
        ('Status: Positive result, methodology-clean (analytical, no seed variance). '
         'Siap dimasukkan ke paper/tesis Section Discussion sebagai empirical justification fusion strategy.',
         8.5, True),
    ], color=RGBColor(0x1F, 0x3A, 0x5F))

    prs.save(PPTX)
    print(f'After: {len(prs.slides)} slides (inserted Space Alignment at 1-idx 132)')
    print(f'Saved: {PPTX}')


if __name__ == '__main__':
    main()
