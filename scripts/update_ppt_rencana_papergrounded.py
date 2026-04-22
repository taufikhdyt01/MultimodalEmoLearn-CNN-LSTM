"""
Restructure PPT slide 128 — "Rencana Eksplorasi Lanjutan"
Dari 8-item list (arahan dosen + turunan paper) → paper-grounded view
supaya dosen paham metodologi eksperimen via paper referensi.

Struktur baru:
  [0] title: Rencana Eksperimen Lanjutan (Paper-Grounded)
  [1] motivasi
  [2] kol kiri: Liliana 2019 family (4 eksperimen — A DONE, B/C/D planning)
  [3] kol kanan: Pitaloka 2017 + Selvaraju 2017
  [4] summary: kenapa paper-driven
  [5] status

Input:  docs/PPT Bimbingan.pptx
Output: in-place (backup ke _pre_papergrounded.pptx)
"""
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

PPTX = Path('docs/PPT Bimbingan.pptx')
BACKUP = Path('docs/PPT Bimbingan_pre_papergrounded.pptx')


def set_text_preserve(text_frame, lines):
    """Replace text_frame content with (text, size_pt, bold) tuples."""
    for para in list(text_frame.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    first_para = text_frame.paragraphs[0]
    for run in list(first_para.runs):
        run._r.getparent().remove(run._r)
    first_para.text = ''
    for i, (txt, size_pt, bold) in enumerate(lines):
        p = first_para if i == 0 else text_frame.add_paragraph()
        run = p.add_run()
        run.text = txt
        if size_pt is not None:
            run.font.size = Pt(size_pt)
        if bold is not None:
            run.font.bold = bold


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    shutil.copy2(PPTX, BACKUP)
    print(f'Backup: {BACKUP}')

    prs = Presentation(PPTX)
    # Slide 128 (1-idx) = index 127 (0-idx) after earlier insert
    slide = prs.slides[127]
    shapes = list(slide.shapes)
    print(f'Slide 128 — {len(shapes)} shapes detected')

    # [0] Title
    set_text_preserve(shapes[0].text_frame, [
        ('Rencana Eksperimen Lanjutan — Paper-Grounded Experiments', 18, True)
    ])

    # [1] Motivasi (anchor eksperimen ke paper referensi)
    set_text_preserve(shapes[1].text_frame, [
        ('Best saat ini: Intermediate TL 4-class B3 = Macro F1 0.521 (val-tuned proper).', 10, True),
        ('Tiap eksperimen dirancang berbasis paper referensi — metodologi jelas, reproducible, dan langsung terhubung ke literatur.', 9.5, False),
    ])

    # [2] Kolom kiri — Liliana 2019 family (4 eksperimen, direct extension paper dosen)
    set_text_preserve(shapes[2].text_frame, [
        ('Liliana et al. (2019) — Cognitive Processing (Springer)', 11, True),
        ('"Fuzzy emotion: a natural approach to automatic facial expression recognition"  [first author = dosen pembimbing]', 8.5, False),
        ('', 3, False),
        ('(A) Soft Label Training  ✅ DONE', 9.5, True),
        ('   → Face API distribusi 7-dim sebagai soft target (fuzzy emotion concept)', 9, False),
        ('   → nb 71 (CNN TL): KL-div 0.517 vs Hard 0.427 (+0.090)', 9, False),
        ('   → nb 72 (Late Fusion TL): 0.437 — gagal transfer ke weighted softmax', 9, False),
        ('   → nb 78 (Intermediate TL): queued — joint training + soft target', 9, False),
        ('', 3, False),
        ('(B) Geometric Features (Table 3 paper)', 9.5, True),
        ('   → 10 komponen facial × 2 metrik (eccentricity + dist ratio) = 20-d GF', 9, False),
        ('   → FCNN_geom (20-d) vs FCNN raw (136-d) — test interpretability', 9, False),
        ('', 3, False),
        ('(C) Geometric + Soft Label Combined  ★★★★★', 9.5, True),
        ('   → Double-extension: 20-d GF + fuzzy soft target (kombinasi A + B)', 9, False),
        ('   → Novelty tertinggi — direct follow-up dari paper dosen', 9, False),
        ('', 3, False),
        ('(D) FEIS Fuzzy Rule-Based (replikasi)', 9.5, True),
        ('   → Mamdani inference, 6 emotion engines — non-DL baseline', 9, False),
        ('   → Bandingkan vs DL: apakah rule-based lebih robust di natural data?', 9, False),
    ])

    # [3] Kolom kanan — Pitaloka 2017 + Selvaraju 2017
    set_text_preserve(shapes[3].text_frame, [
        ('Pitaloka et al. (2017)', 11, True),
        ('"Enhancing CNN with Preprocessing Stage in Automatic Emotion Recognition"', 8.5, False),
        ('', 3, False),
        ('(E) GCN Preprocessing Ablation', 9.5, True),
        ('   → Paper menunjukkan Global Contrast Normalization > min-max', 9, False),
        ('   → Ablation: {min-max, GCN, HistEq} × CNN TL best baseline', 9, False),
        ('   → Quick win: 0.5-1 hari, terukur langsung vs baseline 0.521', 9, False),
        ('', 6, False),
        ('Selvaraju et al. (2017) — Grad-CAM', 11, True),
        ('"Visual Explanations from Deep Networks via Gradient-based Localization"', 8.5, False),
        ('', 3, False),
        ('(F) GradCAM Visualization  🔄 IN-PROGRESS', 9.5, True),
        ('   → Qualitative analysis: fokus model per kelas (eye/mouth/brow)', 9, False),
        ('   → nb 73 scaffolded; run di VPS → docs/figures/gradcam/*.png', 9, False),
        ('   → Supporting figure untuk BAB Discussion tesis', 9, False),
    ])

    # [4] Summary prioritas
    set_text_preserve(shapes[4].text_frame, [
        ('Prioritas berdasarkan paper-grounding + novelty:', 10, True),
        ('', 3, False),
        ('• Liliana family (A-D): direct extension paper dosen → novelty tertinggi untuk tesis + potensi paper follow-up', 9.5, False),
        ('• Pitaloka (E): ablation terukur, quick win — memperkuat preprocessing choice di paper JITeCS', 9.5, False),
        ('• Selvaraju (F): visual explanation — complementary untuk BAB Discussion, tidak push SOTA', 9.5, False),
    ])

    # [5] Status banner
    set_text_preserve(shapes[5].text_frame, [
        ('Status per Apr 2026: A DONE (nb 71/72) + nb 78 queued  |  F in-progress (nb 73)  |  B/C/D/E belum dimulai  |  Detail: docs/eksplorasi_lanjutan.md', 9, True)
    ])

    prs.save(PPTX)
    print(f'Saved: {PPTX}')


if __name__ == '__main__':
    main()
