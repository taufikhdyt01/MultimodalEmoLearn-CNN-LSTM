"""
Delete slide 121 (redundant dengan detail per-dataset di slide 127-130)
+ tambah Early Fusion rows ke Primer Skema 1 di slide 126.

Post-delete slide order:
  121: (was 122) CK+ Skema 1
  122: (was 123) JAFFE Skema 1
  ... dan seterusnya (shift up by 1)
"""
import shutil
from copy import deepcopy
from pathlib import Path
from pptx import Presentation

PPTX   = Path('docs/PPT Bimbingan.pptx')
BACKUP = Path('docs/PPT Bimbingan_pre_delete.pptx')


def set_cell_text_preserve(cell, new_text):
    tf = cell.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        first_run = tf.paragraphs[0].runs[0]
        first_run.text = new_text
        for run in list(tf.paragraphs[0].runs)[1:]:
            run._r.getparent().remove(run._r)
        for para in list(tf.paragraphs)[1:]:
            para._p.getparent().remove(para._p)
    else:
        cell.text = new_text


def clone_and_append_row(table, cells_texts):
    tbl = table._tbl
    new_tr = deepcopy(tbl.tr_lst[-1])
    tbl.append(new_tr)
    new_row = table.rows[len(table.rows) - 1]
    for i, txt in enumerate(cells_texts):
        if i < len(new_row.cells):
            set_cell_text_preserve(new_row.cells[i], txt)


def delete_slide(prs, idx):
    """Remove slide at 0-based idx from sldIdLst (and drop rels)."""
    xml_slides = prs.slides._sldIdLst
    slide_id_elements = list(xml_slides)
    target = slide_id_elements[idx]
    # Remove slide relationship from presentation part
    rId = target.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    # Remove slide ID from list
    xml_slides.remove(target)


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    shutil.copy2(PPTX, BACKUP)
    print(f'Backup: {BACKUP}')

    prs = Presentation(PPTX)
    n_before = len(prs.slides)
    print(f'Before: {n_before} slides')

    # (B) Add EF + EF TL rows to slide 126 Primer Skema 1 (do this FIRST
    # before deleting slide 121, so index 125 is still valid)
    print('\n[1/2] Add EF + EF TL rows to slide 126 Primer Skema 1')
    slide_126 = prs.slides[125]
    tables_126 = [sh.table for sh in slide_126.shapes if sh.has_table]
    primer_7c_add = [
        ['Early Fusion',    '0.246', '0.794', '0.786', '0.794'],
        ['Early Fusion TL', '0.253', '0.713', '0.722', '0.713'],
    ]
    primer_4c_add = [
        ['Early Fusion',    '0.457', '0.822', '0.816', '0.822'],
        ['Early Fusion TL', '0.471', '0.770', '0.770', '0.770'],
    ]
    for cells in primer_7c_add:
        clone_and_append_row(tables_126[0], cells)
        print(f'   7c +row: {cells[0]}')
    for cells in primer_4c_add:
        clone_and_append_row(tables_126[1], cells)
        print(f'   4c +row: {cells[0]}')

    # (A) Delete slide 121 (redundant dengan slide 127-130 per-dataset detail)
    print('\n[2/2] Delete slide 121 (redundant summary)')
    delete_slide(prs, 120)  # 0-based idx 120 = 1-based slide 121

    prs.save(PPTX)
    n_after = len(prs.slides)
    print(f'\nAfter: {n_after} slides (deleted 1, added rows only)')
    print(f'Saved: {PPTX}')


if __name__ == '__main__':
    main()
