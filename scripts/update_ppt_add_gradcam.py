"""
Add SLIDE 34 GradCAM Observasi Awal setelah SLIDE 33 Soft Label.

Position: sekarang SLIDE 33 di index 130 (1-idx 131). GradCAM jadi
slide baru di 1-idx 132, pushing Rencana/Diskusi/Terimakasih turun.

Content:
  - Title + description
  - 2 GradCAM overlay PNG side-by-side (CNN TL baseline vs Intermediate TL B3)
  - Observational notes + keterbatasan metodologi
"""
import shutil
from copy import deepcopy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

PPTX   = Path('docs/PPT Bimbingan.pptx')
BACKUP = Path('docs/PPT Bimbingan_pre_gradcam.pptx')

GRADCAM_DIR = Path('docs/figures/gradcam')
CNN_PNG = GRADCAM_DIR / 'cnn_tl_4c_b1_baseline.png'
INT_PNG = GRADCAM_DIR / 'intermediate_tl_4c_b3_img_branch.png'


def duplicate_slide_after(prs, src_idx):
    src = prs.slides[src_idx]
    new_slide = prs.slides.add_slide(src.slide_layout)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    xml_slides.remove(slides_list[-1])
    xml_slides.insert(src_idx + 1, slides_list[-1])
    return new_slide


def add_textbox(slide, left, top, width, height, text, *,
                size_pt=10, bold=False, italic=False, color=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return tb


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    for p in [CNN_PNG, INT_PNG]:
        if not p.exists():
            raise FileNotFoundError(p)

    shutil.copy2(PPTX, BACKUP)
    print(f'Backup: {BACKUP}')

    prs = Presentation(PPTX)
    # SLIDE 33 Soft Label is at 1-idx 131 → 0-idx 130
    # Insert new slide at 0-idx 131 (so it becomes 1-idx 132)
    src = prs.slides[130]  # use Soft Label slide layout as template
    new_slide = prs.slides.add_slide(src.slide_layout)
    # Remove layout placeholders
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    # Move from end to position 131
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    xml_slides.remove(slides_list[-1])
    xml_slides.insert(131, slides_list[-1])

    # ── Populate content ──
    # Title
    add_textbox(new_slide, 0.3, 0.10, 9.4, 0.40,
                'SLIDE 34: GradCAM Observasi Awal (Eksplorasi, nb 73)',
                size_pt=18, bold=True)
    # Description
    add_textbox(new_slide, 0.3, 0.52, 9.4, 0.28,
                'Visualisasi region citra paling berpengaruh per prediksi kelas — '
                'bandingkan 2 model di 8 sampel (2 per kelas × 4 kelas).',
                size_pt=9.5, italic=True,
                color=RGBColor(0x55, 0x55, 0x55))

    # ── Image 1: CNN TL baseline ──
    new_slide.shapes.add_picture(
        str(CNN_PNG), Inches(0.25), Inches(0.90),
        width=Inches(4.60), height=Inches(2.38),
    )
    add_textbox(new_slide, 0.25, 3.28, 4.60, 0.25,
                'CNN TL 4c B1 (single-modal, Macro F1 = 0.456)',
                size_pt=9, bold=True)

    # ── Image 2: Intermediate TL image branch ──
    new_slide.shapes.add_picture(
        str(INT_PNG), Inches(5.15), Inches(0.90),
        width=Inches(4.60), height=Inches(2.38),
    )
    add_textbox(new_slide, 5.15, 3.28, 4.60, 0.25,
                'Intermediate TL 4c B3 image branch (juara val-tuned, Macro F1 = 0.521)',
                size_pt=9, bold=True)

    # ── Observasi ──
    add_textbox(new_slide, 0.3, 3.62, 9.4, 0.80,
                'Observasi qualitative: CNN baseline → fokus tersebar ke non-emosi region '
                '(aksesori, tangan, pundak, background). Intermediate TL image branch → fokus konsisten '
                'di mouth-nose region di 8/8 sampel, bahkan saat occlusion (tangan menutup mulut).',
                size_pt=9.5, bold=False)

    # ── Keterbatasan ──
    add_textbox(new_slide, 0.3, 4.45, 9.4, 1.10,
                'Keterbatasan metodologi (penting): (1) hanya 2 sampel/kelas — representativeness rendah; '
                '(2) single seed — variance antar run belum diketahui; (3) qualitative only — belum ada '
                'quantitative measure (mis. mean activation di mouth bounding box ≥50 sampel/kelas); '
                '(4) belum ada baseline sanity check (random-init model). Observasi awal ini belum cukup '
                'untuk klaim general — bukan finding paper/tesis, sekadar eksplorasi.',
                size_pt=8.5, italic=True,
                color=RGBColor(0x55, 0x55, 0x55))

    prs.save(PPTX)
    n_after = len(prs.slides)
    print(f'After: {n_after} slides (inserted GradCAM at 1-idx 132)')
    print(f'Saved: {PPTX}')


if __name__ == '__main__':
    main()
