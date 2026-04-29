"""
Insert 3 slides baru sebelum Rencana:
  SLIDE 37 — 3-Class Full Grid + Geometric Features (nb 79+81+82)
  SLIDE 38 — CBAM Attention Negative Result (nb 80)
  SLIDE 39 — Confidence-Stratified Test Analysis (post-hoc)

Position: setelah slide 135 (SLIDE 36 lanjutan), sebelum Rencana (136).
After: Rencana → 139, Diskusi → 140, Paper Q → 141, Terimakasih → 142.
"""
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

PPTX   = Path('docs/PPT Bimbingan.pptx')
BACKUP = Path('docs/PPT Bimbingan_pre_37_39.pptx')


def add_textbox(slide, left, top, width, height, lines, color=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, (txt, size_pt, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i == 0:
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
            p.text = ''
        run = p.add_run()
        run.text = txt
        if size_pt is not None: run.font.size = Pt(size_pt)
        if bold is not None: run.font.bold = bold
        if color: run.font.color.rgb = color
    return tb


def set_cell(cell, txt, size_pt=9, bold=False):
    cell.text = ''
    p = cell.text_frame.paragraphs[0]
    run = p.add_run(); run.text = txt
    run.font.size = Pt(size_pt); run.font.bold = bold


def insert_slide_after(prs, after_idx):
    src = prs.slides[after_idx]
    new_slide = prs.slides.add_slide(src.slide_layout)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    xml_slides.remove(slides_list[-1])
    xml_slides.insert(after_idx + 1, slides_list[-1])
    return new_slide


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    shutil.copy2(PPTX, BACKUP)
    print(f'Backup: {BACKUP}')

    prs = Presentation(PPTX)
    print(f'Before: {len(prs.slides)} slides')

    # ═════════════ SLIDE 37: 3-Class Full Grid + Geometric ═════════════
    # Insert after slide 135 (0-idx 134) → becomes 1-idx 136
    print('\n[1/3] SLIDE 37: 3-Class Full Grid + Geometric Features')
    s37 = insert_slide_after(prs, 134)

    add_textbox(s37, 0.3, 0.08, 9.4, 0.40, [
        ('SLIDE 37: 3-Class Full Grid + Geometric Features (nb 79+81+82)', 17, True)
    ])
    add_textbox(s37, 0.3, 0.50, 9.4, 0.30, [
        ('Pasca nb 82, 27-config grid 3-class lengkap (mirror 7-class earlier). Top-5 by val didominasi Late Fusion.',
         9.5, False),
    ], color=RGBColor(0x55, 0x55, 0x55))

    # Top-5 ranking table (left)
    add_textbox(s37, 0.3, 0.88, 4.6, 0.28, [
        ('Top-5 Ranking (val-based, 27 configs)', 11, True),
    ])
    tbl_shape = s37.shapes.add_table(rows=6, cols=4,
        left=Inches(0.3), top=Inches(1.18), width=Inches(4.6), height=Inches(1.85))
    t = tbl_shape.table
    rows_top5 = [
        ['Rank', 'Config', 'Val F1', 'Test F1'],
        ['1 ⭐', 'Late Fusion TL B3', '0.6229', '0.6370'],
        ['2',   'Late Fusion TL B1', '0.6093', '0.6526'],
        ['3',   'Late Fusion scratch B3', '0.6085', '0.5393'],
        ['4',   'Late Fusion scratch B1', '0.6065', '0.6713'],
        ['5',   'Late Fusion scratch B2', '0.6000', '0.5876'],
    ]
    for ri, row in enumerate(rows_top5):
        for ci, val in enumerate(row):
            set_cell(t.cell(ri, ci), val, size_pt=9, bold=(ri == 0 or ri == 1))

    # Geometric Features mini-table (right)
    add_textbox(s37, 5.10, 0.88, 4.6, 0.28, [
        ('Geometric Features (Liliana 2019, nb 81)', 11, True),
    ])
    add_textbox(s37, 5.10, 1.18, 4.6, 0.30, [
        ('untuk tesis, bukan paper JITeCS', 9, True),
    ], color=RGBColor(0xA8, 0x71, 0x43))
    geo_tbl = s37.shapes.add_table(rows=6, cols=3,
        left=Inches(5.10), top=Inches(1.50), width=Inches(4.6), height=Inches(1.55))
    t2 = geo_tbl.table
    rows_geo = [
        ['Setup', 'Val F1', 'Test F1'],
        ['FCNN_Geometric B1 (20-d)', '0.422', '0.468'],
        ['FCNN_Geometric B3 (20-d)', '0.539', '0.607'],
        ['FCNN_Combined B1 (156-d)', '0.536', '0.592'],
        ['FCNN_Combined B3 (156-d)', '0.561', '0.615'],
        ['LF TL + Combined B3', '0.583', '0.660'],
    ]
    for ri, row in enumerate(rows_geo):
        for ci, val in enumerate(row):
            set_cell(t2.cell(ri, ci), val, size_pt=8.5, bold=(ri == 0))

    add_textbox(s37, 5.10, 3.20, 4.6, 0.45, [
        ('Verdict: gagal beat plain LF TL B3 val=0.6229. 20-d GF lossy vs 136-d raw. '
         'Tidak masuk paper, log di BAB tesis.',
         8.5, False),
    ], color=RGBColor(0x55, 0x55, 0x55))

    # Findings T53-T56 (bottom)
    add_textbox(s37, 0.3, 3.20, 4.6, 2.30, [
        ('Temuan T53-T56 (Konsolidasi)', 10, True),
        ('', 3, False),
        ('T53: Late Fusion strategi paling robust di 3-class (top-5 by val all LF)', 9, True),
        ('   Decision-level averaging optimal untuk coarse class granularity.', 8.5, False),
        ('', 3, False),
        ('T54: TL > scratch konsisten (LF TL B3 0.623 vs scratch 0.609)', 9, True),
        ('   Pre-trained ImageNet boost tetap signifikan di ResNet-18.', 8.5, False),
        ('', 3, False),
        ('T55: Geometric (Liliana) negative — 20-d lossy', 9, True),
        ('   Possibly raw FCNN sudah belajar implicit geometric repr.', 8.5, False),
        ('', 3, False),
        ('T56: w_best shifts dengan branch input dim (0.15 → 0.70)', 9, True),
    ])

    # Status arahan dosen
    add_textbox(s37, 0.3, 5.05, 9.4, 0.45, [
        ('Status arahan dosen: 4/4 DONE  (1) GradCAM ✅  (2) Space Alignment ✅ positive  '
         '(3) CBAM ✅ negative (slide 38)  (4) Geometric ✅ negative (tesis-only)',
         9, True),
    ], color=RGBColor(0x1F, 0x3A, 0x5F))

    # ═════════════ SLIDE 38: CBAM Negative ═════════════
    print('\n[2/3] SLIDE 38: CBAM Attention Negative Result')
    s38 = insert_slide_after(prs, 135)

    add_textbox(s38, 0.3, 0.08, 9.4, 0.40, [
        ('SLIDE 38: CBAM Attention — Negative Result (arahan dosen #3, nb 80)', 17, True)
    ])
    add_textbox(s38, 0.3, 0.50, 9.4, 0.30, [
        ('Test apakah Channel + Spatial Attention (Woo et al. ECCV 2018) bisa boost image stream '
         'untuk beat 3-class juara Late Fusion TL B3 val=0.6229.', 9.5, False),
    ], color=RGBColor(0x55, 0x55, 0x55))

    # Results table
    add_textbox(s38, 0.3, 0.90, 9.4, 0.28, [
        ('Hasil 4 Configs (val-based)', 11, True),
    ])
    cbam_tbl = s38.shapes.add_table(rows=5, cols=6,
        left=Inches(0.3), top=Inches(1.20), width=Inches(9.4), height=Inches(1.45))
    t3 = cbam_tbl.table
    rows_cbam = [
        ['Config', 'Val F1', 'Test F1', 'Test Acc', 'w', 'Δ vs plain baseline'],
        ['CNN_TL_CBAM B1', '0.411', '0.702', '0.833', '—', '−0.082 val (plain 0.493)'],
        ['CNN_TL_CBAM B3', '0.509', '0.649', '0.758', '—', '+0.014 val (plain 0.495) marginal'],
        ['Late Fusion TL CBAM B1', '0.600', '0.652', '0.825', '0.05', '−0.009 val (plain 0.609)'],
        ['Late Fusion TL CBAM B3', '0.609', '0.605', '0.747', '0.00', '−0.014 val (plain 0.623 ⭐)'],
    ]
    for ri, row in enumerate(rows_cbam):
        for ci, val in enumerate(row):
            set_cell(t3.cell(ri, ci), val, size_pt=9, bold=(ri == 0))

    # Key findings
    add_textbox(s38, 0.3, 2.78, 9.4, 1.85, [
        ('Temuan Kunci:', 10, True),
        ('', 3, False),
        ('• CBAM tidak beat baseline plain di semua config — best CBAM val 0.609 < plain LF TL B3 val 0.623', 9.5, False),
        ('• CNN_TL_CBAM B1 val drop −0.082 — attention merusak learning di setup B1 paling sederhana', 9.5, False),
        ('• w_best Late Fusion drop dramatis: 0.25 → 0.05 (B1), 0.15 → 0.00 (B3)', 9.5, True),
        ('     Val grid pilih "ignore CBAM CNN branch entirely" → model konfirmasi CBAM CNN tidak berguna', 9, False),
        ('• Root cause hypothesis: dataset kecil (6,795 train) + natural noise → attention learn spurious features', 9.5, False),
        ('• Konsisten dengan GradCAM (nb 73): CNN baseline fokus ke non-emotion region — attention amplify masalah', 9.5, False),
    ])

    # Decision
    add_textbox(s38, 0.3, 4.70, 9.4, 0.85, [
        ('Decision: stop attention eksplorasi', 11, True),
        ('Tidak extend ke Ghost Module / Triplet Attention. Document sebagai negative finding di tesis.', 9.5, False),
        ('Implikasi paper: Late Fusion TL B3 plain (val 0.6229, test 0.637) tetap juara — confirmed robust.', 9, False),
    ], color=RGBColor(0x8E, 0x44, 0x44))

    # ═════════════ SLIDE 39: Confidence-Stratified Test Analysis ═════════════
    print('\n[3/3] SLIDE 39: Confidence-Stratified Test Analysis')
    s39 = insert_slide_after(prs, 136)

    add_textbox(s39, 0.3, 0.08, 9.4, 0.40, [
        ('SLIDE 39: Confidence-Stratified Test Analysis (post-hoc, no retrain)', 17, True)
    ])
    add_textbox(s39, 0.3, 0.50, 9.4, 0.32, [
        ('Tujuan: separate label-noise effect dari fundamental model limitation. '
         'Evaluasi LF TL B3 di test subset stratified by Face API confidence.',
         9.5, False),
    ], color=RGBColor(0x55, 0x55, 0x55))

    # Results table (left)
    add_textbox(s39, 0.3, 0.90, 4.4, 0.28, [
        ('Hasil per Confidence Threshold', 11, True),
    ])
    conf_tbl = s39.shapes.add_table(rows=7, cols=4,
        left=Inches(0.3), top=Inches(1.20), width=Inches(4.4), height=Inches(2.0))
    t4 = conf_tbl.table
    rows_conf = [
        ['Threshold', 'N', 'Macro F1', 'Catatan'],
        ['≥0.60 (full)', '929', '0.637', 'baseline'],
        ['≥0.70', '892', '0.642', '+0.005'],
        ['≥0.80', '846', '0.639', 'flat'],
        ['≥0.90', '785', '0.645', '+0.008 peak'],
        ['≥0.95', '724', '0.619', 'turun ⬇'],
        ['≥0.99', '631', '0.581', 'turun signifikan ⬇⬇'],
    ]
    for ri, row in enumerate(rows_conf):
        for ci, val in enumerate(row):
            set_cell(t4.cell(ri, ci), val, size_pt=9, bold=(ri == 0))

    # Embed figure (right)
    fig_path = Path('docs/figures/test_macro_f1_by_confidence.png')
    if fig_path.exists():
        s39.shapes.add_picture(str(fig_path), Inches(5.00), Inches(1.20),
                                width=Inches(4.70), height=Inches(2.20))

    # Interpretation
    add_textbox(s39, 0.3, 3.40, 9.4, 1.20, [
        ('Interpretasi: Refutes Label-Noise Hypothesis', 11, True),
        ('', 3, False),
        ('• Macro F1 flat 0.62-0.65 (th 0.60-0.90) lalu DROP di th ≥0.95 — bukan monotonic naik seperti yang diharapkan', 9.5, False),
        ('• Negative class F1 collapse 0.328 → 0.071 di conf≥0.99 (sample minority hilang di high-conf subset)', 9.5, False),
        ('• Akar masalah: distribution shift (Face API confident di neutral, minority biasanya conf 0.4-0.7)', 9.5, False),
    ])

    # Conclusion
    add_textbox(s39, 0.3, 4.70, 9.4, 0.85, [
        ('Kesimpulan: model limitation = minority class learning under extreme imbalance, BUKAN label noise', 10, True),
        ('Validasi conf60 sebagai optimal trade-off (preserve sample sufficiency + minority class viability).', 9, False),
        ('Untuk paper Discussion: refute hipotesis label-noise-limited, frame minority class sebagai core challenge.', 9, False),
    ], color=RGBColor(0x1F, 0x3A, 0x5F))

    prs.save(PPTX)
    print(f'\nAfter: {len(prs.slides)} slides (inserted 3 new slides)')


if __name__ == '__main__':
    main()
