"""
Split SLIDE 32 pair-dataset slides menjadi per-dataset (1 dataset per slide).

Before (4 source slides, 4 tables each = overflow):
  122: CK+ & JAFFE (Skema 1)
  123: RAF-DB & KDEF (Skema 1)
  125: CK+ & JAFFE → Primer (Skema 2)
  126: RAF-DB & KDEF → Primer (Skema 2)

After (8 slides, 2 tables each — fits in single screen):
  122: CK+ (Skema 1)           124: RAF-DB (Skema 1)
  123: JAFFE (Skema 1)         125: KDEF (Skema 1)
  126: Primer (unchanged)
  127: CK+ → Primer            129: RAF-DB → Primer
  128: JAFFE → Primer          130: KDEF → Primer
"""
import shutil
from copy import deepcopy
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

PPTX   = Path('docs/PPT Bimbingan.pptx')
BACKUP = Path('docs/PPT Bimbingan_pre_split.pptx')


def set_text_first_run(text_frame, new_text):
    """Replace text_frame content, preserve first run's style."""
    for para in list(text_frame.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    first_para = text_frame.paragraphs[0]
    if first_para.runs:
        first_run = first_para.runs[0]
        first_run.text = new_text
        for run in list(first_para.runs)[1:]:
            run._r.getparent().remove(run._r)
    else:
        first_para.text = new_text


def duplicate_slide_after(prs, src_idx):
    """Clone slide at src_idx, insert at src_idx+1. Return new slide."""
    src = prs.slides[src_idx]
    new_slide = prs.slides.add_slide(src.slide_layout)
    # Remove layout placeholder shapes
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    # Deep-copy shapes from source
    for shape in src.shapes:
        new_slide.shapes._spTree.append(deepcopy(shape._element))
    # Move to position src_idx + 1
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    xml_slides.remove(slides_list[-1])
    xml_slides.insert(src_idx + 1, slides_list[-1])
    return new_slide


def delete_shapes_by_idx(slide, indices_to_delete):
    """Delete shapes at given indices (0-based) from slide."""
    shapes = list(slide.shapes)
    # Sort descending so indices remain valid while deleting
    for idx in sorted(indices_to_delete, reverse=True):
        if idx < len(shapes):
            sp = shapes[idx]._element
            sp.getparent().remove(sp)


def shift_shapes_y(slide, indices, dy_inches):
    """Shift Y position of shapes by dy_inches (negative = up)."""
    shapes = list(slide.shapes)
    delta_emu = int(dy_inches * 914400)
    for idx in indices:
        if idx < len(shapes):
            sh = shapes[idx]
            sh.top = Emu(sh.top + delta_emu)


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    shutil.copy2(PPTX, BACKUP)
    print(f'Backup: {BACKUP}')

    prs = Presentation(PPTX)
    n_before = len(prs.slides)

    # Process in DESCENDING order so earlier indices don't shift
    # Slide configs (0-idx based):
    #   Slide 122 (idx 121): CK+ & JAFFE (Skema 1)
    #   Slide 123 (idx 122): RAF-DB & KDEF (Skema 1)
    #   Slide 125 (idx 124): CK+ & JAFFE → Primer (Skema 2)
    #   Slide 126 (idx 125): RAF-DB & KDEF → Primer (Skema 2)

    # ─────── Slide 126 (idx 125): RAF-DB & KDEF → Primer ───────
    print('\n[1/4] Split slide 126 (RAF-DB & KDEF → Primer)')
    src = prs.slides[125]
    # Shapes: 0=bg, 1=title, 2=RAF7c label, 3=RAF7c tbl, 4=RAF4c label, 5=RAF4c tbl,
    #         6=KDEF7c label, 7=KDEF7c tbl, 8=KDEF4c label, 9=KDEF4c tbl, 10=conclusion
    new = duplicate_slide_after(prs, 125)
    # Original = RAF-DB only: delete KDEF tables + conclusion
    delete_shapes_by_idx(src, [6, 7, 8, 9, 10])
    set_text_first_run(list(src.shapes)[1].text_frame,
                       'SLIDE 32 (lanjutan): Skema 2 — RAF-DB → Primer')
    # New = KDEF only: delete RAF-DB tables, shift KDEF shapes up, keep conclusion
    delete_shapes_by_idx(new, [2, 3, 4, 5])
    # After deletion, indices in 'new': 0=bg, 1=title, 2=KDEF7c label, 3=KDEF7c tbl,
    # 4=KDEF4c label, 5=KDEF4c tbl, 6=conclusion
    shift_shapes_y(new, [2, 3, 4, 5], dy_inches=-2.50)
    set_text_first_run(list(new.shapes)[1].text_frame,
                       'SLIDE 32 (lanjutan): Skema 2 — KDEF → Primer')

    # ─────── Slide 125 (idx 124): CK+ & JAFFE → Primer ───────
    print('[2/4] Split slide 125 (CK+ & JAFFE → Primer)')
    src = prs.slides[124]
    # Shapes: 0=bg, 1=title, 2=intro, 3=CK7c label, 4=CK7c tbl, 5=CK4c label, 6=CK4c tbl,
    #         7=JAFFE7c label, 8=JAFFE7c tbl, 9=JAFFE4c label, 10=JAFFE4c tbl, 11=conclusion
    new = duplicate_slide_after(prs, 124)
    # Original = CK+ only
    delete_shapes_by_idx(src, [7, 8, 9, 10, 11])
    set_text_first_run(list(src.shapes)[1].text_frame,
                       'SLIDE 32 (lanjutan): Skema 2 — CK+ → Primer')
    # New = JAFFE only
    delete_shapes_by_idx(new, [3, 4, 5, 6])
    # After deletion: 0=bg, 1=title, 2=intro, 3=JAFFE7c label, 4=JAFFE7c tbl,
    # 5=JAFFE4c label, 6=JAFFE4c tbl, 7=conclusion
    shift_shapes_y(new, [3, 4, 5, 6], dy_inches=-2.50)
    set_text_first_run(list(new.shapes)[1].text_frame,
                       'SLIDE 32 (lanjutan): Skema 2 — JAFFE → Primer')

    # ─────── Slide 123 (idx 122): RAF-DB & KDEF ───────
    print('[3/4] Split slide 123 (RAF-DB & KDEF Skema 1)')
    src = prs.slides[122]
    # Shapes: 0=bg, 1=title, 2=RAF7c label, 3=RAF7c tbl, 4=RAF4c label, 5=RAF4c tbl,
    #         6=KDEF7c label, 7=KDEF7c tbl, 8=KDEF4c label, 9=KDEF4c tbl, 10=conclusion
    new = duplicate_slide_after(prs, 122)
    delete_shapes_by_idx(src, [6, 7, 8, 9, 10])
    set_text_first_run(list(src.shapes)[1].text_frame,
                       'SLIDE 32 (lanjutan): Skema 1 Lengkap — RAF-DB')
    delete_shapes_by_idx(new, [2, 3, 4, 5])
    shift_shapes_y(new, [2, 3, 4, 5], dy_inches=-2.84)
    set_text_first_run(list(new.shapes)[1].text_frame,
                       'SLIDE 32 (lanjutan): Skema 1 Lengkap — KDEF')

    # ─────── Slide 122 (idx 121): CK+ & JAFFE ───────
    print('[4/4] Split slide 122 (CK+ & JAFFE Skema 1)')
    src = prs.slides[121]
    # Shapes: 0=bg, 1=title, 2=instruction, 3=CK7c label, 4=CK7c tbl, 5=CK4c label, 6=CK4c tbl,
    #         7=JAFFE7c label, 8=JAFFE7c tbl, 9=JAFFE4c label, 10=JAFFE4c tbl
    new = duplicate_slide_after(prs, 121)
    delete_shapes_by_idx(src, [7, 8, 9, 10])
    set_text_first_run(list(src.shapes)[1].text_frame,
                       'SLIDE 32: Skema 1 Lengkap — CK+ (+ Late Fusion TL)')
    delete_shapes_by_idx(new, [3, 4, 5, 6])
    # After deletion: 0=bg, 1=title, 2=instruction, 3=JAFFE7c label, 4=JAFFE7c tbl,
    # 5=JAFFE4c label, 6=JAFFE4c tbl
    shift_shapes_y(new, [3, 4, 5, 6], dy_inches=-2.84)
    set_text_first_run(list(new.shapes)[1].text_frame,
                       'SLIDE 32 (lanjutan): Skema 1 Lengkap — JAFFE')

    prs.save(PPTX)
    n_after = len(prs.slides)
    print(f'\nBefore: {n_before} slides  →  After: {n_after} slides')
    print(f'Saved: {PPTX}')


if __name__ == '__main__':
    main()
