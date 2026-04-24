"""
Drop 4-class references dari slide 133 & 135 karena paper reframe ke 3-class.
Keep 7-class sebagai fine-grained baseline saja.
"""
import shutil
from copy import deepcopy
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

PPTX   = Path('docs/PPT Bimbingan.pptx')
BACKUP = Path('docs/PPT Bimbingan_pre_drop4c.pptx')


def set_text_preserve(tf, lines):
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    first_para = tf.paragraphs[0]
    for run in list(first_para.runs):
        run._r.getparent().remove(run._r)
    first_para.text = ''
    for i, (txt, size_pt, bold) in enumerate(lines):
        p = first_para if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = txt
        if size_pt is not None: run.font.size = Pt(size_pt)
        if bold is not None: run.font.bold = bold


def set_cell(cell, txt, size_pt=9, bold=False):
    cell.text = ''
    p = cell.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(size_pt)
    run.font.bold = bold


def remove_table_row(table, ridx):
    tbl = table._tbl
    rows = tbl.tr_lst
    if 0 <= ridx < len(rows):
        tbl.remove(rows[ridx])


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    shutil.copy2(PPTX, BACKUP)
    print(f'Backup: {BACKUP}')

    prs = Presentation(PPTX)

    # ═════ Slide 133: fix comparison table + description + findings ═════
    s133 = prs.slides[132]
    shapes = list(s133.shapes)
    print('\n[Slide 133]')

    # Update description (shape 1) — remove "4-class mapping arbitrary" phrasing
    set_text_preserve(shapes[1].text_frame, [
        ('Mapping 3-class valence (Russell 1980 circumplex): positive=happy+surprised, '
         'neutral, negative=sad+angry+fearful+disgusted. Imbalance 1:14 — jauh lebih '
         'balanced dari 7-class (1:1138). 5 arch × 3 scenario = 15 configs, val-based selection.',
         9.5, False),
    ])
    print('  shape[1] description: updated')

    # Find the 4-row comparison table (Scheme | Juara | Test F1)
    tables = [sh for sh in shapes if sh.has_table]
    comp_tbl = None
    for sh in tables:
        if len(sh.table.rows) == 4 and len(sh.table.columns) == 3:
            comp_tbl = sh.table
            break
    if comp_tbl is not None:
        # Remove row with "4-class" — it's row index 2 based on inspection
        remove_table_row(comp_tbl, 2)
        # Now table has 3 rows: header, 7c, 3c
        print('  comparison table: dropped 4-class row')
    else:
        print('  [WARN] comparison table not found')

    # Update findings text (shape 4) — remove "4c" references in Temuan
    set_text_preserve(shapes[4].text_frame, [
        ('Temuan Kunci (T49-T52)', 10, True),
        ('', 3, False),
        ('T49: 3-class valence mapping punya literature precedent (Russell 1980)', 9, True),
        ('Imbalance 1:14 vs 1:1138 di 7-class — jauh lebih balanced untuk evaluasi minoritas.',
         8.5, False),
        ('', 4, False),
        ('T50: Late Fusion TL B3 juara val-based (0.623)', 9, True),
        ('Decision-level averaging optimal untuk coarse class granularity.', 8.5, False),
        ('', 4, False),
        ('T51: w_best = 0.15-0.30 (lebih balanced)', 9, True),
        ('Di 3-class, image stream contribute signifikan (tidak FCNN-dominant seperti 7-class).',
         8.5, False),
        ('', 4, False),
        ('T52: FCNN 3c sangat kompetitif (test=0.634)', 9, True),
        ('Landmark geometry dominan — 136-d coord vector cukup untuk valence discrimination.',
         8.5, False),
    ])
    print('  shape[4] findings: updated (removed 4c mentions)')

    # ═════ Slide 135 Rencana: update motivasi + Eksplorasi Lain ═════
    s135 = prs.slides[134]
    shapes135 = list(s135.shapes)
    print('\n[Slide 135]')

    # shape[1] motivasi
    set_text_preserve(shapes135[1].text_frame, [
        ('Best saat ini: 3-class Late Fusion TL B3 = Macro F1 0.623 (val-tuned, nb 79).', 10, True),
        ('Status 4 arahan dosen + turunan paper. Arahan diprioritaskan.', 9.5, False),
    ])
    print('  shape[1] motivasi: removed "reframe dari 4c"')

    # shape[3] Eksplorasi Lain — remove 4c comparison
    set_text_preserve(shapes135[3].text_frame, [
        ('Turunan + Eksplorasi Lain', 11, True),
        ('', 3, False),
        ('Soft Label (nb 71/72/78) ✅ CLOSED — negative result', 9, True),
        ('', 3, False),
        ('3-Class Exploration (nb 79) ✅ DONE — positive', 9, True),
        ('   → Late Fusion TL B3 val-based 0.623', 8.5, False),
        ('', 3, False),
        ('Pitaloka 2017 GCN Ablation ⏳ belum', 9, True),
        ('   → quick win 0.5-1 hari', 8.5, False),
        ('', 3, False),
        ('Liliana 2019 FEIS Fuzzy Rule ⏳ belum', 9, True),
        ('   → non-DL baseline, est. 2-3 hari', 8.5, False),
    ])
    print('  shape[3] Eksplorasi Lain: removed "+0.10 vs 4c"')

    prs.save(PPTX)
    print(f'\nSaved: {PPTX}')


if __name__ == '__main__':
    main()
