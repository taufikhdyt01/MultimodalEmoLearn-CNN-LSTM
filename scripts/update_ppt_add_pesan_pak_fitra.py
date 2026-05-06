"""
Add slide "Update Komunikasi ke Pak Fitra — Submission JITeCS" before Terimakasih.

Isi: pesan WA ke Pak Fitra — update draft JITeCS v0.4 (reframe 7c+3c, sudah
review Pak Budi) + pertanyaan timeline submission (target sebelumnya 24 April).
"""
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

PPTX   = Path('docs/PPT Bimbingan.pptx')
BACKUP = Path('docs/PPT Bimbingan_pre_pesan_fitra.pptx')


def add_textbox(slide, left, top, width, height, lines, *, color=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, (txt, size_pt, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i == 0:
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
            p.text = ''
        run = p.add_run()
        run.text = txt
        if size_pt is not None:
            run.font.size = Pt(size_pt)
        if bold is not None:
            run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return tb


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    shutil.copy2(PPTX, BACKUP)
    print(f'Backup: {BACKUP}')

    prs = Presentation(PPTX)
    n_before = len(prs.slides)
    print(f'Before: {n_before} slides')

    # Terimakasih saat ini di posisi terakhir (0-idx n-1).
    # Pertanyaan Bimbingan paper (slide 141, 0-idx 140) adalah anchor.
    # Insert new slide tepat sebelum Terimakasih.
    last_idx = n_before - 1  # Terimakasih
    src = prs.slides[last_idx]
    new_slide = prs.slides.add_slide(src.slide_layout)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    # Move new slide ke posisi sebelum Terimakasih
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    new_sld = slides_list[-1]
    xml_slides.remove(new_sld)
    xml_slides.insert(last_idx, new_sld)

    # ── Title ──
    add_textbox(new_slide, 0.3, 0.10, 9.4, 0.40,
                [('Update Komunikasi ke Pak Fitra — Submission JITeCS', 18, True)])

    # ── Pesan WA (verbatim, sopan) ──
    add_textbox(new_slide, 0.3, 0.60, 9.4, 0.35,
                [('Pesan terkirim ke Pak Fitra (update draft + pertanyaan timeline)',
                  11, True)],
                color=RGBColor(0x2E, 0x5E, 0x96))

    # Greeting block
    add_textbox(new_slide, 0.5, 1.00, 9.0, 0.55,
                [('Selamat siang, Pak Fitra.', 10.5, False),
                 ('Mohon maaf mengganggu waktunya.', 10.5, False)])

    # Body 1 — update draft v0.4
    add_textbox(new_slide, 0.5, 1.65, 9.0, 0.95,
                [('Saat ini draft paper untuk submission ke JITeCS sudah saya revisi menjadi versi v0.4. '
                  'Ada perubahan pada hasil yang dilaporkan, dari sebelumnya 7 class & 4 class menjadi '
                  '7 class & 3 class yang sudah disesuaikan dengan literatur. Selain itu, draft ini juga '
                  'sudah mendapatkan review dan masukan dari Pak Budi.', 10.5, False)])

    # Drive link block
    add_textbox(new_slide, 0.5, 2.75, 9.0, 0.32,
                [('Untuk file terbaru bisa diakses di drive yang sama berikut pak:', 10.5, False)])
    add_textbox(new_slide, 0.5, 3.07, 9.0, 0.32,
                [('https://drive.google.com/drive/folders/1tbxb0Vi1AsEuXPxNNkRqECmX24P9Kqfj?usp=drive_link',
                  9, False)],
                color=RGBColor(0x2E, 0x5E, 0x96))

    # Body 2 — pertanyaan timeline submission
    add_textbox(new_slide, 0.5, 3.55, 9.0, 1.10,
                [('Saya juga ingin menanyakan terkait timeline submission ke JITeCS, karena sebelumnya '
                  'pada timeline akselerasi yang Bapak sampaikan targetnya tanggal 24 April. '
                  'Kira-kira untuk draft ini akan dikirim kapan ya Pak? Mohon arahan dari Bapak.',
                  10.5, False)])

    # Closing
    add_textbox(new_slide, 0.5, 4.85, 9.0, 0.40,
                [('Terima kasih banyak Pak.', 10.5, True)],
                color=RGBColor(0x55, 0x55, 0x55))

    # Footer note
    add_textbox(new_slide, 0.3, 6.85, 9.4, 0.30,
                [('Status: menunggu balasan Pak Fitra terkait jadwal submission.',
                  9, False)],
                color=RGBColor(0x88, 0x88, 0x88))

    prs.save(PPTX)
    print(f'After:  {len(prs.slides)} slides (+1 = "Update Komunikasi ke Pak Fitra")')
    print(f'Saved:  {PPTX}')


if __name__ == '__main__':
    main()
