"""
Update PPT Bimbingan — add missing rows + val-tuned fixes for slides 122-126.

Changes:
  Slide 122 (Skema 1 — CK+/JAFFE):
    - Update Late Fusion row (val-tuned): CK+ 7c/4c, JAFFE 7c/4c
    - Add Early Fusion + Early Fusion TL rows (from nb 66)

  Slide 123 (Skema 1 — RAF-DB/KDEF):
    - Add Early Fusion + Early Fusion TL rows (from nb 66)

  Slide 125 (Skema 2 — CK+/JAFFE → Primer):
    - Add Early Fusion + Early Fusion TL rows (nb 68)
    - Add Late Fusion TL row (nb 69)

  Slide 126 (Skema 2 — RAF-DB/KDEF → Primer):
    - Add Early Fusion + Early Fusion TL rows (nb 68)
    - Add Late Fusion TL row (nb 69)

Input:  docs/PPT Bimbingan.pptx
Output: docs/PPT Bimbingan.pptx (in-place)

Backup akan disimpan ke docs/PPT Bimbingan_pre_ef_lf.pptx sebelum write.
"""
import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation

PPTX = Path('docs/PPT Bimbingan.pptx')
BACKUP = Path('docs/PPT Bimbingan_pre_ef_lf.pptx')

# ───────── Data ─────────

# Slide 122: CK+ 7c, CK+ 4c, JAFFE 7c, JAFFE 4c (4 tables, order matches slide layout)
SLIDE_122 = [
    # CK+ 7c
    {
        'update_row': (6, ['Late Fusion', '0.494', '0.780', '0.691', '0.780']),
        'add_rows': [
            ['Early Fusion',    '0.446', '0.695', '0.665', '0.695'],
            ['Early Fusion TL', '0.762', '0.847', '0.847', '0.847'],
        ],
    },
    # CK+ 4c
    {
        'update_row': (6, ['Late Fusion', '0.537', '0.694', '0.681', '0.694']),
        'add_rows': [
            ['Early Fusion',    '0.507', '0.694', '0.667', '0.694'],
            ['Early Fusion TL', '0.795', '0.871', '0.872', '0.871'],
        ],
    },
    # JAFFE 7c
    {
        'update_row': (6, ['Late Fusion', '0.314', '0.400', '0.290', '0.400']),
        'add_rows': [
            ['Early Fusion',    '0.286', '0.350', '0.267', '0.350'],
            ['Early Fusion TL', '0.041', '0.150', '0.043', '0.150'],
        ],
    },
    # JAFFE 4c
    {
        'update_row': (6, ['Late Fusion', '0.492', '0.650', '0.615', '0.650']),
        'add_rows': [
            ['Early Fusion',    '0.177', '0.550', '0.390', '0.550'],
            ['Early Fusion TL', '0.352', '0.600', '0.507', '0.600'],
        ],
    },
]

# Slide 123: RAF-DB 7c, RAF-DB 4c, KDEF 7c, KDEF 4c
SLIDE_123 = [
    # RAF-DB 7c
    {'add_rows': [
        ['Early Fusion',    '0.710', '0.808', '0.804', '0.808'],
        ['Early Fusion TL', '0.693', '0.790', '0.786', '0.790'],
    ]},
    # RAF-DB 4c
    {'add_rows': [
        ['Early Fusion',    '0.792', '0.818', '0.819', '0.818'],
        ['Early Fusion TL', '0.799', '0.823', '0.823', '0.823'],
    ]},
    # KDEF 7c
    {'add_rows': [
        ['Early Fusion',    '0.667', '0.674', '0.663', '0.674'],
        ['Early Fusion TL', '0.799', '0.795', '0.797', '0.795'],
    ]},
    # KDEF 4c
    {'add_rows': [
        ['Early Fusion',    '0.693', '0.763', '0.764', '0.763'],
        ['Early Fusion TL', '0.816', '0.855', '0.854', '0.855'],
    ]},
]

# Slide 125: Skema 2 CK+/JAFFE → Primer (current 7 rows incl header; need add 3)
SLIDE_125 = [
    # CK+ → Primer 7c
    {'add_rows': [
        ['Early Fusion',    '0.125', '0.686', '0.638', '0.686'],
        ['Early Fusion TL', '0.179', '0.705', '0.676', '0.705'],
        ['Late Fusion TL',  '0.247', '0.703', '0.726', '0.703'],
    ]},
    # CK+ → Primer 4c
    {'add_rows': [
        ['Early Fusion',    '0.207', '0.577', '0.582', '0.577'],
        ['Early Fusion TL', '0.101', '0.194', '0.288', '0.194'],
        ['Late Fusion TL',  '0.141', '0.182', '0.281', '0.182'],
    ]},
    # JAFFE → Primer 7c
    {'add_rows': [
        ['Early Fusion',    '0.026', '0.032', '0.026', '0.032'],
        ['Early Fusion TL', '0.001', '0.002', '0.000', '0.002'],
        ['Late Fusion TL',  '0.050', '0.118', '0.192', '0.118'],
    ]},
    # JAFFE → Primer 4c
    {'add_rows': [
        ['Early Fusion',    '0.004', '0.009', '0.000', '0.009'],
        ['Early Fusion TL', '0.004', '0.009', '0.000', '0.009'],
        ['Late Fusion TL',  '0.015', '0.013', '0.009', '0.013'],
    ]},
]

# Slide 126: Skema 2 RAF-DB/KDEF → Primer
SLIDE_126 = [
    # RAF-DB → Primer 7c
    {'add_rows': [
        ['Early Fusion',    '0.142', '0.365', '0.476', '0.365'],
        ['Early Fusion TL', '0.157', '0.685', '0.663', '0.685'],
        ['Late Fusion TL',  '0.153', '0.483', '0.575', '0.483'],
    ]},
    # RAF-DB → Primer 4c
    {'add_rows': [
        ['Early Fusion',    '0.227', '0.404', '0.485', '0.404'],
        ['Early Fusion TL', '0.311', '0.477', '0.575', '0.477'],
        ['Late Fusion TL',  '0.212', '0.459', '0.517', '0.459'],
    ]},
    # KDEF → Primer 7c
    {'add_rows': [
        ['Early Fusion',    '0.007', '0.010', '0.011', '0.010'],
        ['Early Fusion TL', '0.029', '0.054', '0.097', '0.054'],
        ['Late Fusion TL',  '0.055', '0.078', '0.063', '0.078'],
    ]},
    # KDEF → Primer 4c
    {'add_rows': [
        ['Early Fusion',    '0.039', '0.040', '0.023', '0.040'],
        ['Early Fusion TL', '0.104', '0.101', '0.152', '0.101'],
        ['Late Fusion TL',  '0.017', '0.011', '0.003', '0.011'],
    ]},
]


# ───────── Helpers ─────────

def set_cell_text_preserve(cell, new_text):
    """Replace cell text while keeping first run's font/style."""
    tf = cell.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        first_run = tf.paragraphs[0].runs[0]
        first_run.text = new_text
        # Remove extra runs in first para
        for run in list(tf.paragraphs[0].runs)[1:]:
            run._r.getparent().remove(run._r)
        # Remove extra paragraphs
        for para in list(tf.paragraphs)[1:]:
            para._p.getparent().remove(para._p)
    else:
        cell.text = new_text


def clone_last_row_with_text(table, cell_texts):
    """Clone last row's XML (preserves formatting), append, set new text."""
    tbl = table._tbl
    last_tr = tbl.tr_lst[-1]
    new_tr = deepcopy(last_tr)
    tbl.append(new_tr)
    new_row = table.rows[len(table.rows) - 1]
    for i, text in enumerate(cell_texts):
        if i < len(new_row.cells):
            set_cell_text_preserve(new_row.cells[i], text)


def get_tables(slide):
    return [sh.table for sh in slide.shapes if sh.has_table]


def apply_update(slide, specs, slide_num):
    tables = get_tables(slide)
    print(f'\n Slide {slide_num}: {len(tables)} tables found, {len(specs)} specs')
    if len(tables) != len(specs):
        print(f'   [WARN] tables({len(tables)}) != specs({len(specs)})')
    for ti, (tbl, spec) in enumerate(zip(tables, specs)):
        # Update row (Late Fusion fix)
        if 'update_row' in spec:
            ridx, texts = spec['update_row']
            row = tbl.rows[ridx]
            for ci, t in enumerate(texts):
                if ci < len(row.cells):
                    set_cell_text_preserve(row.cells[ci], t)
            print(f'   [table {ti}] updated row {ridx}: {texts[0]} → {texts[1]}')
        # Add rows
        for cells in spec.get('add_rows', []):
            clone_last_row_with_text(tbl, cells)
            print(f'   [table {ti}] +row: {cells[0]} ({cells[1]})')


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    shutil.copy2(PPTX, BACKUP)
    print(f'Backup: {BACKUP}')

    prs = Presentation(PPTX)
    apply_update(prs.slides[122 - 1], SLIDE_122, 122)
    apply_update(prs.slides[123 - 1], SLIDE_123, 123)
    apply_update(prs.slides[125 - 1], SLIDE_125, 125)
    apply_update(prs.slides[126 - 1], SLIDE_126, 126)

    prs.save(PPTX)
    print(f'\nSaved: {PPTX}')


if __name__ == '__main__':
    main()
