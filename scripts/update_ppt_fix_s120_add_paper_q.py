"""
(A) Fix slide 120 — Late Fusion row di table 5 fusion strategy salah isi
    (tercampur dengan Early/Intermediate values). Update Temuan 32 & 35
    dengan angka val-tuned proper.

(B) Tambah slide baru "Pertanyaan Bimbingan — Paper Updates" sebelum
    Terimakasih: EECCIS 2026 (accepted) + JITeCS (draft 14 halaman).
"""
import shutil
from copy import deepcopy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

PPTX   = Path('docs/PPT Bimbingan.pptx')
BACKUP = Path('docs/PPT Bimbingan_pre_paperq.pptx')


def set_cell_text_preserve(cell, new_text):
    tf = cell.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        first_run = tf.paragraphs[0].runs[0]
        first_run.text = new_text
        for run in list(tf.paragraphs[0].runs)[1:]:
            run._r.getparent().remove(run._r)
        for para in list(tf.paragraphs)[1:]:
            para._p.getparent().remove(para._p)
    else:
        cell.text = new_text


def set_text_preserve(text_frame, lines):
    """Replace TF content with [(txt, size, bold), ...]."""
    for para in list(text_frame.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    first_para = text_frame.paragraphs[0]
    for run in list(first_para.runs):
        run._r.getparent().remove(run._r)
    first_para.text = ''
    for i, (txt, size_pt, bold) in enumerate(lines):
        p = first_para if i == 0 else text_frame.add_paragraph()
        run = p.add_run()
        run.text = txt
        if size_pt is not None:
            run.font.size = Pt(size_pt)
        if bold is not None:
            run.font.bold = bold


def add_textbox(slide, left, top, width, height, lines, *, color=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, (txt, size_pt, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # Clear default runs in first para
        if i == 0:
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
            p.text = ''
        run = p.add_run()
        run.text = txt
        if size_pt is not None:
            run.font.size = Pt(size_pt)
        if bold is not None:
            run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return tb


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    shutil.copy2(PPTX, BACKUP)
    print(f'Backup: {BACKUP}')

    prs = Presentation(PPTX)
    print(f'Before: {len(prs.slides)} slides')

    # ─────── (A) Fix slide 120 ───────
    print('\n[A] Fix slide 120 (Fusion Strategy table + Temuan)')
    slide_120 = prs.slides[119]
    shapes_120 = list(slide_120.shapes)

    # Fix table row values + stars (val-tuned proper)
    tbl = shapes_120[3].table
    # r3 Early Fusion — add ★ on 7c (juara 7-class overall 0.333)
    set_cell_text_preserve(tbl.rows[3].cells[1], '0.333 (TL B3) ★')
    # r4 Intermediate — add ★ on 4c (juara 4-class overall 0.521)
    set_cell_text_preserve(tbl.rows[4].cells[2], '0.521 (TL B3) ★')
    # r5 Late Fusion — fix values (val-tuned: 7c best = 0.270 scratch B1, 4c best = 0.479 scratch B2)
    set_cell_text_preserve(tbl.rows[5].cells[1], '0.270 (scratch B1)')
    set_cell_text_preserve(tbl.rows[5].cells[2], '0.479 (scratch B2)')
    print('   table: Late Fusion row fixed + stars repositioned')

    # Update Temuan 32 (old 0.301 test-tuned → 0.270 val-tuned)
    set_text_preserve(shapes_120[5].text_frame, [
        ('Temuan 32: Early Fusion TL B3 tembus best di 7-class (0.333)', 9.5, True),
        ('Melampaui Intermediate TL B3 (0.292) dan Late Fusion scratch B1 (0.270, val-tuned). Kombinasi heatmap channel + class weights + augmentasi + TL bekerja sinergis untuk skenario kelas banyak + imbalance ekstrem.', 9, False),
    ])
    print('   Temuan 32: updated Late Fusion number (0.301 → 0.270 val-tuned)')

    # Update Temuan 35 (ranking 4-class salah post val-tuning)
    set_text_preserve(shapes_120[8].text_frame, [
        ('Temuan 35: Ranking fusion strategy berbeda per granularitas kelas (val-tuned)', 9.5, True),
        ('7-class: Early TL B3 (0.333) > Intermediate TL B3 (0.292) > Late Fusion scratch B1 (0.270). 4-class: Intermediate TL B3 (0.521) > Late Fusion scratch B2 (0.479) > Early TL B1 (0.471). Early Fusion juara di high-granularity, Intermediate juara di low-granularity.', 9, False),
    ])
    print('   Temuan 35: updated ranking val-tuned')

    # ─────── (B) Add "Pertanyaan Paper" slide before Terimakasih ───────
    print('\n[B] Add new slide "Pertanyaan Bimbingan — Paper Updates"')
    # Terimakasih is currently at 1-idx 134 → 0-idx 133
    # Use slide 133 layout as template
    src = prs.slides[133]
    new_slide = prs.slides.add_slide(src.slide_layout)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    # Move to position 133 (before Terimakasih, which is at 133)
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    xml_slides.remove(slides_list[-1])
    xml_slides.insert(133, slides_list[-1])

    # Content:
    # Title
    add_textbox(new_slide, 0.3, 0.10, 9.4, 0.40,
                [('Pertanyaan Bimbingan — Update Paper', 18, True)])

    # ── Section 1: EECCIS 2026 ──
    add_textbox(new_slide, 0.3, 0.60, 9.4, 0.35,
                [('1. Paper EECCIS 2026 (Accepted)', 13, True)],
                color=RGBColor(0x2E, 0x5E, 0x96))
    add_textbox(new_slide, 0.3, 0.98, 9.4, 0.30,
                [('Judul: "Automated Engagement Assessment in Block-Based Programming Using Facial Emotion Recognition"', 10, False)])
    add_textbox(new_slide, 0.3, 1.30, 9.4, 1.10,
                [('Pertanyaan:', 10, True),
                 ('• Langkah-langkah revisi final manuscript — apa yang perlu diperhatikan dari reviewer comments?', 9.5, False),
                 ('• Proses registrasi konferensi — deadline, author registration, early-bird fee?', 9.5, False),
                 ('• Mekanisme pengajuan pendanaan untuk biaya konferensi — apakah ada skema dari fakultas/universitas?', 9.5, False)])

    # ── Section 2: JITeCS ──
    add_textbox(new_slide, 0.3, 2.55, 9.4, 0.35,
                [('2. Paper JITeCS (Draft)', 13, True)],
                color=RGBColor(0xA8, 0x71, 0x43))
    add_textbox(new_slide, 0.3, 2.93, 9.4, 0.30,
                [('Status: Draft selesai — tersimpan di folder Publikasi + Google Drive template', 10, False)])
    add_textbox(new_slide, 0.3, 3.25, 9.4, 0.30,
                [('https://drive.google.com/drive/folders/1tbxb0Vi1AsEuXPxNNkRqECmX24P9Kqfj', 8.5, False)],
                color=RGBColor(0x2E, 0x5E, 0x96))
    add_textbox(new_slide, 0.3, 3.60, 9.4, 1.15,
                [('Kendala + pertanyaan:', 10, True),
                 ('• Draft saat ini 14 halaman — melebihi 2 halaman dari batas (12 halaman).', 9.5, False),
                 ('• Saya belum menemukan bagian yang sebaiknya dipangkas — mohon masukan + review.', 9.5, False),
                 ('• Apakah ada sesi bimbingan hari ini untuk membahas revisi bersama?', 9.5, False)])

    # Footer closing
    add_textbox(new_slide, 0.3, 4.95, 9.4, 0.50,
                [('Terima kasih banyak atas bimbingannya, Pak.', 10, True)],
                color=RGBColor(0x55, 0x55, 0x55))

    prs.save(PPTX)
    print(f'\nAfter: {len(prs.slides)} slides')
    print(f'Saved: {PPTX}')


if __name__ == '__main__':
    main()
