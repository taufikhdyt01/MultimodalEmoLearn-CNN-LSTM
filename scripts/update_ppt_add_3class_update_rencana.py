"""
Update PPT:
  (A) Insert 2 slides SLIDE 36 (3-Class Exploration) setelah SLIDE 35 Space Alignment
  (B) Update slide Rencana Eksperimen Lanjutan dengan current status

After run:
  133: SLIDE 36 3-Class Results (new)
  134: SLIDE 36 lanjutan 3-Class CM + Best Model (new)
  135: Rencana Eksplorasi (was 133, updated status)
  136: Diskusi (was 134)
  137: Pertanyaan Bimbingan Paper (was 135)
  138: Terimakasih (was 136)
"""
import shutil
from copy import deepcopy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

PPTX   = Path('docs/PPT Bimbingan.pptx')
BACKUP = Path('docs/PPT Bimbingan_pre_3class.pptx')

CM_3C_PNG = Path('docs/figures/confusion_matrix_3class.png')


def set_text_preserve(tf, lines):
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    first_para = tf.paragraphs[0]
    for run in list(first_para.runs):
        run._r.getparent().remove(run._r)
    first_para.text = ''
    for i, (txt, size_pt, bold) in enumerate(lines):
        p = first_para if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = txt
        if size_pt is not None: run.font.size = Pt(size_pt)
        if bold is not None: run.font.bold = bold


def add_textbox(slide, left, top, width, height, lines, color=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, (txt, size_pt, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i == 0:
            for r in list(p.runs):
                r._r.getparent().remove(r._r)
            p.text = ''
        run = p.add_run()
        run.text = txt
        if size_pt is not None: run.font.size = Pt(size_pt)
        if bold is not None: run.font.bold = bold
        if color: run.font.color.rgb = color
    return tb


def insert_blank_slide_at(prs, src_idx_for_layout, target_pos):
    src = prs.slides[src_idx_for_layout]
    new_slide = prs.slides.add_slide(src.slide_layout)
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    xml_slides.remove(slides_list[-1])
    xml_slides.insert(target_pos, slides_list[-1])
    return new_slide


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    if not CM_3C_PNG.exists():
        raise FileNotFoundError(CM_3C_PNG)

    shutil.copy2(PPTX, BACKUP)
    print(f'Backup: {BACKUP}')

    prs = Presentation(PPTX)
    print(f'Before: {len(prs.slides)} slides')

    # ═══════════════════════════════════════════════════
    # (A) Insert 2 slides for SLIDE 36 at position 132 & 133 (0-idx)
    # ═══════════════════════════════════════════════════
    # SLIDE 35 Space Alignment is at 0-idx 131 (1-idx 132)
    # Insert SLIDE 36 results at 0-idx 132 (becomes 1-idx 133)
    # Insert SLIDE 36 lanjutan at 0-idx 133 (becomes 1-idx 134)

    # ── Slide A: SLIDE 36 Results ──
    new_A = insert_blank_slide_at(prs, src_idx_for_layout=131, target_pos=132)

    add_textbox(new_A, 0.3, 0.08, 9.4, 0.40, [
        ('SLIDE 36: 3-Class Exploration (Positive/Neutral/Negative, nb 79)', 18, True)
    ])
    add_textbox(new_A, 0.3, 0.50, 9.4, 0.30, [
        ('4-class mapping arbitrary → 3-class valence (Russell 1980 circumplex). Imbalance 1:14 '
         '(4× lebih balanced dari 4c). 5 arch × 3 scenario = 15 configs, val-based selection.',
         9.5, False),
    ], color=RGBColor(0x55, 0x55, 0x55))

    # ── 15-config table (compact) ──
    tbl_shape = new_A.shapes.add_table(
        rows=16, cols=5,
        left=Inches(0.3), top=Inches(0.88), width=Inches(5.3), height=Inches(3.8))
    tbl = tbl_shape.table

    header = ['Config', 'Val F1', 'Test F1', 'Test Acc', 'w']
    rows_data = [
        ['FCNN B1',              '0.603', '0.589', '0.741', '—'],
        ['FCNN B2',              '0.564', '0.575', '0.702', '—'],
        ['FCNN B3',              '0.619', '0.634', '0.749', '—'],
        ['CNN TL B1',            '0.493', '0.634', '0.791', '—'],
        ['CNN TL B2',            '0.515', '0.516', '0.581', '—'],
        ['CNN TL B3',            '0.495', '0.706', '0.840', '—'],
        ['Intermediate TL B1',   '0.554', '0.686', '0.790', '—'],
        ['Intermediate TL B2',   '0.559', '0.665', '0.815', '—'],
        ['Intermediate TL B3',   '0.501', '0.689', '0.828', '—'],
        ['Late Fusion TL B1',    '0.609', '0.653', '0.797', '0.25'],
        ['Late Fusion TL B2',    '0.591', '0.646', '0.774', '0.30'],
        ['Late Fusion TL B3 ⭐', '0.623', '0.637', '0.784', '0.15'],
        ['Early Fusion TL B1',   '0.537', '0.587', '0.783', '—'],
        ['Early Fusion TL B2',   '0.504', '0.584', '0.650', '—'],
        ['Early Fusion TL B3',   '0.509', '0.699', '0.820', '—'],
    ]
    for ci, h in enumerate(header):
        cell = tbl.cell(0, ci)
        cell.text = ''
        p = cell.text_frame.paragraphs[0]
        run = p.add_run(); run.text = h
        run.font.size = Pt(8); run.font.bold = True
    for ri, row in enumerate(rows_data, start=1):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]
            run = p.add_run(); run.text = val
            run.font.size = Pt(8)
            if '⭐' in row[0] or '**' in val:
                run.font.bold = True

    # ── Right side: Comparison + Findings ──
    add_textbox(new_A, 5.80, 0.88, 3.90, 0.30, [
        ('Perbandingan vs scheme lain (val-based)', 10, True),
    ])

    comp_tbl = new_A.shapes.add_table(
        rows=4, cols=3,
        left=Inches(5.80), top=Inches(1.18), width=Inches(3.90), height=Inches(1.1))
    t = comp_tbl.table
    comp_data = [
        ['Scheme', 'Juara', 'Test F1'],
        ['7-class', 'Early Fusion TL B3', '0.333'],
        ['4-class', 'Intermediate TL B3', '0.521'],
        ['3-class ⭐', 'Late Fusion TL B3', '0.637'],
    ]
    for ri, row in enumerate(comp_data):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]
            run = p.add_run(); run.text = val
            run.font.size = Pt(9); run.font.bold = (ri == 0 or '⭐' in row[0])

    add_textbox(new_A, 5.80, 2.40, 3.90, 3.10, [
        ('Temuan Kunci (T49-T52)', 10, True),
        ('', 3, False),
        ('T49: Gain +0.10 Macro F1 dari 4c → 3c', 9, True),
        ('Semua arch lebih stabil, bahkan config terlemah CNN TL B2 test=0.516 masih di level 4c best.',
         8.5, False),
        ('', 4, False),
        ('T50: Fusion strategy shifts per class granularity', 9, True),
        ('7c: Early Fusion → 4c: Intermediate → 3c: Late Fusion. Late Fusion reclaims juara di coarse.',
         8.5, False),
        ('', 4, False),
        ('T51: w_best lebih balanced (0.15-0.30)', 9, True),
        ('vs 4-class (0.00-0.15 FCNN-dominant). Image stream contribute signifikan di 3-class.',
         8.5, False),
        ('', 4, False),
        ('T52: FCNN 3c sangat kompetitif', 9, True),
        ('FCNN B3 val=0.619 nyaris setara Late Fusion TL — landmark geometry dominan di Primer.',
         8.5, False),
    ])

    # ── Slide B: SLIDE 36 lanjutan — Best Model + CM ──
    new_B = insert_blank_slide_at(prs, src_idx_for_layout=131, target_pos=133)

    add_textbox(new_B, 0.3, 0.08, 9.4, 0.40, [
        ('SLIDE 36 (lanjutan): 3-Class Best Model — Late Fusion TL B3', 18, True)
    ])
    add_textbox(new_B, 0.3, 0.50, 9.4, 0.28, [
        ('Val Macro = 0.623  |  Test Macro = 0.637  |  Test Acc = 0.784  |  w_best = 0.15',
         10, True),
    ], color=RGBColor(0x1F, 0x3A, 0x5F))

    # ── Per-class metrics table (left) ──
    add_textbox(new_B, 0.3, 0.88, 4.3, 0.28, [
        ('Per-Class Metrics', 10, True),
    ])
    pc_tbl = new_B.shapes.add_table(
        rows=6, cols=5,
        left=Inches(0.3), top=Inches(1.18), width=Inches(4.3), height=Inches(1.8))
    t = pc_tbl.table
    pc_data = [
        ['Class', 'Prec', 'Recall', 'F1', 'Support'],
        ['positive', '0.607', '0.930', '0.735', '186'],
        ['neutral',  '0.935', '0.776', '0.848', '688'],
        ['negative', '0.288', '0.382', '0.328', '55'],
        ['Macro avg', '0.610', '0.696', '0.637', '929'],
        ['Weighted avg', '0.844', '0.784', '0.795', '929'],
    ]
    for ri, row in enumerate(pc_data):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]
            run = p.add_run(); run.text = val
            run.font.size = Pt(9)
            run.font.bold = (ri == 0 or ri >= 4)

    # ── Confusion matrix image (right) ──
    new_B.shapes.add_picture(str(CM_3C_PNG),
                              Inches(5.00), Inches(0.90),
                              width=Inches(4.80), height=Inches(3.55))

    # ── Misclassification summary bottom ──
    add_textbox(new_B, 0.3, 3.20, 4.3, 2.35, [
        ('Misclassification Pattern', 10, True),
        ('', 3, False),
        ('• neutral → positive: 105 (15.3%) — dominant, subtle smile reactions', 8.5, False),
        ('• negative → neutral: 27 (49.1% of negative) — frustrasi mild ≈ deep focus', 8.5, False),
        ('• neutral → negative: 49 (7.1%)', 8.5, False),
        ('• Cross-valence confusion total: 13 (1.4%) — rendah', 8.5, False),
        ('', 4, False),
        ('Strong valence discrimination:', 9, True),
        ('positive ↔ negative confusion hampir nol. Konsisten Russell 1980 — valence dimension paling prominent.',
         8.5, False),
        ('Negative recall rendah (0.382) karena boundary neutral-negative ambiguous di natural data.',
         8.5, False),
    ])

    # ═══════════════════════════════════════════════════
    # (B) Update Rencana slide (now at 0-idx 134, was 132 pre-insert)
    # ═══════════════════════════════════════════════════
    # After inserts: Rencana slide moved from 1-idx 133 → 1-idx 135 (0-idx 134)
    rencana = prs.slides[134]
    shapes = list(rencana.shapes)

    # Update title
    set_text_preserve(shapes[0].text_frame, [
        ('Rencana & Status Eksplorasi Lanjutan (Fokus Arahan Dosen)', 18, True)
    ])

    # Motivasi — update best claim
    set_text_preserve(shapes[1].text_frame, [
        ('Best saat ini: 3-class Late Fusion TL B3 = Macro F1 0.623 (val-tuned, nb 79) — reframe dari 4c.', 10, True),
        ('Status 4 arahan dosen + turunan paper. Arahan diprioritaskan.', 9.5, False),
    ])

    # Kolom kiri — arahan dosen (4 items)
    set_text_preserve(shapes[2].text_frame, [
        ('Arahan Dosen (4 item)', 11, True),
        ('', 3, False),
        ('(1) GradCAM ✅ DONE (nb 73) — observational', 9.5, True),
        ('(2) Space Alignment ✅ DONE (nb 75)', 9.5, True),
        ('   → positive result: CCA top-5 = 0.978, siap paper', 8.5, False),
        ('(3) Attention Module (CBAM) 🔄 nb 80 prepared', 9.5, True),
        ('   → target beat 3c LF TL B3 (0.623), ~4-5 jam VPS', 8.5, False),
        ('(4) Geometric Features (Liliana 2019) ⏳ belum', 9.5, True),
        ('   → 20-d GF Table 3, est. 2-3 hari', 8.5, False),
    ])

    # Kolom kanan — turunan paper + 3-class positive
    set_text_preserve(shapes[3].text_frame, [
        ('Turunan + Eksplorasi Lain', 11, True),
        ('', 3, False),
        ('Soft Label (nb 71/72/78) ✅ CLOSED — negative result', 9, True),
        ('', 3, False),
        ('3-Class Exploration (nb 79) ✅ DONE — positive', 9, True),
        ('   → +0.10 Macro F1 vs 4c, reframe paper', 8.5, False),
        ('', 3, False),
        ('Pitaloka 2017 GCN Ablation ⏳ belum', 9, True),
        ('   → quick win 0.5-1 hari', 8.5, False),
        ('', 3, False),
        ('Liliana 2019 FEIS Fuzzy Rule ⏳ belum', 9, True),
        ('   → non-DL baseline, est. 2-3 hari', 8.5, False),
    ])

    # Prioritas urutan
    set_text_preserve(shapes[4].text_frame, [
        ('Urutan eksekusi rekomendasi:', 10, True),
        ('', 3, False),
        ('(3) Attention CBAM (nb 80 queued) → (4) Geometric Features → Pitaloka GCN → FEIS. '
         'Semua berbasis val-based selection, bandingkan vs 3-class baseline 0.623.', 9, False),
    ])

    # Status banner
    set_text_preserve(shapes[5].text_frame, [
        ('Status Apr 2026: 3/4 arahan dosen (1, 2) + 3-class DONE  |  (3) queued (nb 80)  |  '
         '(4) + 2 turunan paper belum dimulai  |  Detail: docs/eksplorasi_lanjutan.md, bimbingan_progress.md',
         9, True),
    ])

    prs.save(PPTX)
    n_after = len(prs.slides)
    print(f'After: {n_after} slides (inserted 2 + updated Rencana)')


if __name__ == '__main__':
    main()
