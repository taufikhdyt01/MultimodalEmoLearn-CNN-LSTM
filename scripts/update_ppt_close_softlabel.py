"""
Update PPT slide 130 (SLIDE 33 Soft Label) — close eksplorasi sebagai
negative result. Update title, table (3 arch summary), Temuan 41-44.
"""
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt

PPTX   = Path('docs/PPT Bimbingan.pptx')
BACKUP = Path('docs/PPT Bimbingan_pre_closesoft.pptx')


def set_cell_text_preserve(cell, new_text):
    tf = cell.text_frame
    if tf.paragraphs and tf.paragraphs[0].runs:
        first_run = tf.paragraphs[0].runs[0]
        first_run.text = new_text
        for run in list(tf.paragraphs[0].runs)[1:]:
            run._r.getparent().remove(run._r)
        for para in list(tf.paragraphs)[1:]:
            para._p.getparent().remove(para._p)
    else:
        cell.text = new_text


def set_text_preserve(text_frame, lines):
    """Replace text_frame with (text, size_pt, bold) tuples."""
    for para in list(text_frame.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    first_para = text_frame.paragraphs[0]
    for run in list(first_para.runs):
        run._r.getparent().remove(run._r)
    first_para.text = ''
    for i, (txt, size_pt, bold) in enumerate(lines):
        p = first_para if i == 0 else text_frame.add_paragraph()
        run = p.add_run()
        run.text = txt
        if size_pt is not None:
            run.font.size = Pt(size_pt)
        if bold is not None:
            run.font.bold = bold


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    shutil.copy2(PPTX, BACKUP)
    print(f'Backup: {BACKUP}')

    prs = Presentation(PPTX)
    slide = prs.slides[129]  # SLIDE 33 = 1-idx 130 = 0-idx 129
    shapes = list(slide.shapes)

    # [1] Title
    set_text_preserve(shapes[1].text_frame, [
        ('SLIDE 33: Soft Label Training — Eksplorasi Selesai (Negative Result)', 18, True)
    ])

    # [2] Subtitle
    set_text_preserve(shapes[2].text_frame, [
        ('Konsolidasi nb 71 (CNN TL) + nb 72 (Late Fusion TL) + nb 78 (Intermediate TL) — 4 loss × 3 arsitektur, '
         'hyperparam align (EPOCHS=50, LR=5e-5), selection by val macro.', 10, False)
    ])

    # [3] Table 6×4 — rewrite summary
    tbl = shapes[3].table
    rows_data = [
        ['Arsitektur',            'A Hard CE', 'C KL-div', 'Kesimpulan'],
        ['CNN TL 4c B1 (nb 71)',  '0.432',     '0.427',    'Tied dalam variance (val N/A)'],
        ['Late Fusion TL 4c (nb 72)', '—',     '0.437',    '< baseline val-tuned 0.466'],
        ['Intermediate TL 4c B1 (nb 78)', '0.485', '0.453', 'Hard CE juara by val (0.480)'],
        ['Variance single-seed',  '±0.05-0.09', '—',       'khas DL — multi-seed needed'],
        ['Hipotesis',             'Tidak terkonfirmasi', '—', 'STOP eksplorasi'],
    ]
    for ri, row_cells in enumerate(rows_data):
        for ci, txt in enumerate(row_cells):
            set_cell_text_preserve(tbl.rows[ri].cells[ci], txt)

    # [4] Temuan header
    set_text_preserve(shapes[4].text_frame, [('Temuan Konsolidasi (41-44)', 12, True)])

    # [5-8] Temuan boxes
    temuan = [
        ('T41: Single-seed variance ±0.05-0.09 — terlalu besar untuk klaim marginal',
         'Gap antar loss config 0.02-0.03, di bawah variance. Mustahil klaim "soft label menang" tanpa multi-seed (3-5 run minimum).'),
        ('T42: KL-div juara 0.517 di run sebelumnya adalah fluke',
         'Re-run dengan setup konsisten kasih KL-div = 0.427 di CNN TL (selisih 0.09 dari run lama) — manifestasi pure seed/convergence variance.'),
        ('T43: Soft label tidak transfer ke fusion',
         'Late Fusion TL gagal (0.432-0.437 vs baseline 0.466). Intermediate TL: Hard CE menang by val. Modality interaction dominates over loss-function nuance.'),
        ('T44: Status — stop eksplorasi soft label',
         'Hipotesis tidak terkonfirmasi. Fokus pindah ke Geometric Features (Liliana Table 3) + GCN Preprocessing (Pitaloka 2017) yang potensinya lebih jelas.'),
    ]
    for i, (heading, body) in enumerate(temuan):
        shp = shapes[5 + i]
        set_text_preserve(shp.text_frame, [
            (heading, 9.5, True),
            (body, 9, False),
        ])

    prs.save(PPTX)
    print(f'Saved: {PPTX}')


if __name__ == '__main__':
    main()
