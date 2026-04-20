"""
Update PPT Bimbingan dengan val-tuned Late Fusion numbers (Primer conf60).

Changes:
  - Best overall: "Late Fusion TL 4-class B3 = 0.567" → "Intermediate TL 4-class B3 = 0.521"
  - Best 7c: "Late Fusion TL B1 = 0.301" → "Early Fusion TL B3 = 0.333"
  - Primer Late Fusion + Late Fusion TL rows di SLIDE 32 Primer tables

Slides affected (1-indexed):
  114 = SLIDE 26 Conf60 BREAKTHROUGH
  115 = Temuan Confidence Filtering
  116 = SLIDE 27 Conf60 + Undersampling
  117 = SLIDE 28 Analisis Masalah Dataset (konsultasi)
  120 = SLIDE 30 Early Fusion (lanjutan comparison)
  121 = SLIDE 31 Cross-Dataset
  124 = SLIDE 32c Primer Skema 1 + Summary

Tidak disentuh (context lain):
  slide 99 (cell 0.567 dari dataset lain)
  slide 122-123 (CK+/JAFFE/RAF-DB/KDEF tables — 0.567 di RAF-DB/CK+ bukan Primer)
"""
from pathlib import Path
from pptx import Presentation

PPTX = Path('d:/MultimodalEmoLearn/docs/PPT Bimbingan.pptx')


def replace_in_text_frame(tf, replacements):
    """In-place replace substrings in a text frame across runs."""
    for para in tf.paragraphs:
        for run in para.runs:
            for old, new in replacements.items():
                if old in run.text:
                    run.text = run.text.replace(old, new)


def replace_in_shape(shape, replacements):
    if shape.has_text_frame:
        replace_in_text_frame(shape.text_frame, replacements)
    elif shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                replace_in_text_frame(cell.text_frame, replacements)


def apply_slide(prs, slide_idx_1based, replacements, context=''):
    """Apply replacements to specific slide (1-indexed)."""
    slide = prs.slides[slide_idx_1based - 1]
    print(f'  Slide {slide_idx_1based}{(" ("+context+")") if context else ""}:')
    for shape in slide.shapes:
        replace_in_shape(shape, replacements)


def main():
    prs = Presentation(str(PPTX))

    # ═══ Slide 114: SLIDE 26 Conf60 BREAKTHROUGH ═══
    apply_slide(prs, 114, {
        'Best overall: Late Fusion TL 4-class B3': 'Best overall: Intermediate TL 4-class B3',
        '0.567': '0.521',
        '0.301 (Late Fusion TL B1)': '0.333 (Early Fusion TL B3)',
        '0.301': '0.333',  # for standalone 7c best
    }, 'SLIDE 26 BREAKTHROUGH')

    # ═══ Slide 115: Temuan Confidence Filtering ═══
    apply_slide(prs, 115, {
        '0.412 -> 0.567': '0.412 -> 0.521 (val-tuned proper)',
        '0.412 → 0.567': '0.412 → 0.521 (val-tuned proper)',
        '+38%': '+26%',
        'Late Fusion TL B3': 'Intermediate TL B3',
        '0.567': '0.521',
    }, 'Temuan CF')

    # ═══ Slide 116: SLIDE 27 Conf60 + Undersampling ═══
    apply_slide(prs, 116, {
        'Late Fusion TL 4c B3': 'Intermediate TL 4c B3',
        'Late Fusion TL 4-class B3 conf60 = 0.567': 'Intermediate TL 4-class B3 conf60 = 0.521 (val-tuned)',
        '0.567': '0.521',
    }, 'SLIDE 27 Kombinasi')

    # ═══ Slide 117: SLIDE 28 Konsultasi ═══
    apply_slide(prs, 117, {
        'F1 0.567': 'F1 0.521',
    }, 'SLIDE 28 Konsultasi')

    # ═══ Slide 120: SLIDE 30 Early Fusion lanjutan (comparison table) ═══
    apply_slide(prs, 120, {
        '0.301 (TL B1)': '0.333 (Early TL B3)',    # 7c comparison
        '0.567 (TL B3) ★': '0.521 (Inter TL B3) ★',  # 4c comparison
        '0.567 (TL B3)': '0.521 (Inter TL B3)',
        'melampaui Intermediate TL B3 (0.292) dan Late Fusion TL B1 (0.301)':
            'di atas Intermediate TL B3 (0.292) dan Late Fusion TL B1 (0.238)',
        '0.471 vs 0.567 Late Fusion': '0.471 vs 0.521 Intermediate TL',
    }, 'SLIDE 30 Early Fusion lanjutan')

    # ═══ Slide 121: SLIDE 31 Cross-Dataset ═══
    apply_slide(prs, 121, {
        'Vs Primer self-training best 0.567': 'Vs Primer self-training best 0.521',
        'Primer self 0.567': 'Primer self 0.521',
        '0.567': '0.521',
    }, 'SLIDE 31 Cross-Dataset')

    # ═══ Slide 124: SLIDE 32c Primer tables + summary ═══
    apply_slide(prs, 124, {
        'Late Fusion TL 4c B3 = 0.567': 'Intermediate TL 4c B3 = 0.521',
        'Best Primer keseluruhan (dengan B3 + TL + augmentation) = Late Fusion TL 4c B3 = 0.567':
            'Best Primer keseluruhan (val-tuned w, proper) = Intermediate TL 4c B3 = 0.521',
        'best Primer 4c = 0.482': 'best Primer 4c (val-tuned) = 0.521',
        # Primer Late Fusion val-tuned row updates (7c + 4c)
        # 7c: Late Fusion 0.244→0.270, Late Fusion TL 0.285→0.238
        # 4c: Late Fusion 0.460→0.474, Late Fusion TL 0.472→0.422
    }, 'SLIDE 32c Primer')

    # Primer Late Fusion table rows di slide 124 — update specific cells
    slide = prs.slides[123]  # 0-indexed for slide 124
    for shape in slide.shapes:
        if shape.has_table:
            tbl = shape.table
            # Find rows by first column content
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells]
                if not cells: continue
                first = cells[0]
                # Primer 7c table: Late Fusion row
                if first == 'Late Fusion' and len(cells) >= 5 and cells[1] == '0.244':
                    row.cells[1].text_frame.paragraphs[0].runs[0].text = '0.270'
                    row.cells[2].text_frame.paragraphs[0].runs[0].text = '0.816'
                    row.cells[3].text_frame.paragraphs[0].runs[0].text = '0.812'
                    row.cells[4].text_frame.paragraphs[0].runs[0].text = '0.816'
                # Primer 7c Late Fusion TL row
                elif first == 'Late Fusion TL' and len(cells) >= 5 and cells[1] == '0.285':
                    row.cells[1].text_frame.paragraphs[0].runs[0].text = '0.238'
                    row.cells[2].text_frame.paragraphs[0].runs[0].text = '0.790'
                    row.cells[3].text_frame.paragraphs[0].runs[0].text = '0.784'
                    row.cells[4].text_frame.paragraphs[0].runs[0].text = '0.790'
                # Primer 4c Late Fusion
                elif first == 'Late Fusion' and len(cells) >= 5 and cells[1] == '0.460':
                    row.cells[1].text_frame.paragraphs[0].runs[0].text = '0.474'
                    row.cells[2].text_frame.paragraphs[0].runs[0].text = '0.807'
                    row.cells[3].text_frame.paragraphs[0].runs[0].text = '0.815'
                    row.cells[4].text_frame.paragraphs[0].runs[0].text = '0.807'
                # Primer 4c Late Fusion TL
                elif first == 'Late Fusion TL' and len(cells) >= 5 and cells[1] == '0.472':
                    row.cells[1].text_frame.paragraphs[0].runs[0].text = '0.422'
                    row.cells[2].text_frame.paragraphs[0].runs[0].text = '0.695'
                    row.cells[3].text_frame.paragraphs[0].runs[0].text = '0.722'
                    row.cells[4].text_frame.paragraphs[0].runs[0].text = '0.695'

    prs.save(str(PPTX))
    print(f'\nSaved: {PPTX}')


if __name__ == '__main__':
    main()
