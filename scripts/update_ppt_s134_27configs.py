"""
Patch slide 134 (3-Class Exploration):
- Tabel utama: 16 rows (header + 15 configs) → 28 rows (header + 27 configs)
- Tambah 12 row scratch (CNN/Intermediate/Late Fusion/Early Fusion × B1/B2/B3)
- Macro column → ganti nilai jadi val_macro_f1 (val-based selection criterion)
- Micro/W-F1/Acc/w tetap dari test (val-only metrics tidak di-log per config)
- Update subtitle: "5 arch × 3 scenario = 15 configs" → "9 arch × 3 scenario = 27 configs"
- Add footnote disclaimer di bawah tabel
"""
import json
import shutil
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

PPTX = Path("docs/PPT Bimbingan.pptx")
BACKUP = Path("docs/PPT Bimbingan_pre_s134_27cfg.pptx")
RESULTS_TL = Path("models/frontonly_conf60/3class/all_results_3class.json")
RESULTS_SC = Path("models/frontonly_conf60/3class/scratch_all_results.json")


def w_str(v):
    return f"{v:.2f}" if isinstance(v, (int, float)) and v else "—"


def build_rows():
    """Return list of (label, val_macro, test_macro, test_micro, test_wf1, test_acc, w_or_dash)."""
    tl = json.loads(RESULTS_TL.read_text())
    sc = json.loads(RESULTS_SC.read_text())

    def row(d, key, label):
        v = d[key]
        return (
            label,
            f"{v['val_macro_f1']:.3f}",
            f"{v['test_macro_f1']:.3f}",
            f"{v['test_micro_f1']:.3f}",
            f"{v['test_weighted_f1']:.3f}",
            f"{v['test_accuracy']:.3f}",
            w_str(v.get("best_cnn_weight")),
        )

    # Pair scratch + TL by arch family for readability
    rows = []
    # FCNN (no scratch counterpart — landmark only)
    rows += [row(tl, f"FCNN_B{i}", f"FCNN B{i}") for i in (1, 2, 3)]
    # CNN scratch + TL
    rows += [row(sc, f"CNN_scratch_B{i}", f"CNN scratch B{i}") for i in (1, 2, 3)]
    rows += [row(tl, f"CNN_TL_B{i}", f"CNN TL B{i}") for i in (1, 2, 3)]
    # Intermediate scratch + TL
    rows += [row(sc, f"Intermediate_scratch_B{i}", f"Intermediate scratch B{i}") for i in (1, 2, 3)]
    rows += [row(tl, f"Intermediate_TL_B{i}", f"Intermediate TL B{i}") for i in (1, 2, 3)]
    # Late Fusion scratch + TL  (TL B3 = juara ⭐)
    rows += [row(sc, f"Late_Fusion_scratch_B{i}", f"Late Fusion scratch B{i}") for i in (1, 2, 3)]
    lf_tl = [row(tl, f"Late_Fusion_TL_B{i}", f"Late Fusion TL B{i}") for i in (1, 2, 3)]
    # Mark juara on B3
    lf_tl[2] = ("Late Fusion TL B3 ⭐",) + lf_tl[2][1:]
    rows += lf_tl
    # Early Fusion scratch + TL
    rows += [row(sc, f"Early_Fusion_scratch_B{i}", f"Early Fusion scratch B{i}") for i in (1, 2, 3)]
    rows += [row(tl, f"Early_Fusion_TL_B{i}", f"Early Fusion TL B{i}") for i in (1, 2, 3)]
    return rows


def clear_cell(cell):
    tf = cell.text_frame
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    p = tf.paragraphs[0]
    for run in list(p.runs):
        run._r.getparent().remove(run._r)
    p.text = ""


def set_cell(cell, text, *, size=7, bold=False, color=None, align_center=False):
    clear_cell(cell)
    tf = cell.text_frame
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    if align_center:
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    shutil.copy2(PPTX, BACKUP)
    print(f"Backup: {BACKUP}")

    prs = Presentation(PPTX)
    sl = prs.slides[133]  # slide 134 (1-idx) = 0-idx 133

    rows_data = build_rows()
    assert len(rows_data) == 27, f"expected 27 rows, got {len(rows_data)}"

    # ── (1) Update subtitle text ──
    sub = sl.shapes[1].text_frame
    # Replace "5 arch × 3 scenario = 15 configs" with "9 arch × 3 scenario = 27 configs"
    for para in sub.paragraphs:
        for run in para.runs:
            if run.text and "15 configs" in run.text:
                run.text = run.text.replace("5 arch × 3 scenario = 15 configs",
                                            "9 arch × 3 scenario = 27 configs")
                print("   subtitle: updated 15->27 configs")

    # ── (2) Patch big table (shape index 5) ──
    table_shape = sl.shapes[5]
    tbl = table_shape.table
    n_rows_now = len(tbl.rows)
    n_cols = len(tbl.columns)
    print(f"   table before: {n_rows_now}x{n_cols}")

    # Need 28 rows total (1 header + 27 data). Currently 16 → need to add 12 more.
    # Clone last row 12 times
    tbl_xml = tbl._tbl
    last_tr = list(tbl_xml.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}tr"))[-1]
    for _ in range(12):
        new_tr = deepcopy(last_tr)
        last_tr.addnext(new_tr)
        last_tr = new_tr
    print(f"   table after add: {len(tbl.rows)} rows")

    # Update header row: column 1 label "Macro" → "Val F1"
    header_cells = tbl.rows[0].cells
    set_cell(header_cells[0], "Model", size=7.5, bold=True, align_center=False,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    set_cell(header_cells[1], "Val F1", size=7.5, bold=True, align_center=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    set_cell(header_cells[2], "Macro", size=7.5, bold=True, align_center=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    set_cell(header_cells[3], "Micro", size=7.5, bold=True, align_center=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    set_cell(header_cells[4], "W-F1", size=7.5, bold=True, align_center=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    set_cell(header_cells[5], "Acc", size=7.5, bold=True, align_center=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    # Note: existing table is 6 cols, but we need 7 (Model, Val F1, Macro, Micro, W-F1, Acc, w)
    # Wait — original is 6 cols. Need to insert one more col OR drop one.
    # Decision: KEEP 6 cols. Repurpose existing "Macro" col → Val F1.
    # Map: Model | Val F1 (was Macro) | Micro | W-F1 | Acc | w
    # Reset header properly:
    set_cell(tbl.rows[0].cells[0], "Model", size=7.5, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    set_cell(tbl.rows[0].cells[1], "Val F1", size=7.5, bold=True, align_center=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    set_cell(tbl.rows[0].cells[2], "Micro*", size=7.5, bold=True, align_center=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    set_cell(tbl.rows[0].cells[3], "W-F1*", size=7.5, bold=True, align_center=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    set_cell(tbl.rows[0].cells[4], "Acc*", size=7.5, bold=True, align_center=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    set_cell(tbl.rows[0].cells[5], "w", size=7.5, bold=True, align_center=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))

    # ── (3) Fill 27 data rows (rows 1..27) ──
    juara_color = RGBColor(0xC8, 0x4B, 0x12)  # orange highlight for ⭐
    for i, (label, val_f1, _tmac, tmic, twfi, tacc, ww) in enumerate(rows_data):
        r = tbl.rows[i + 1]
        is_juara = "⭐" in label
        col = juara_color if is_juara else None
        set_cell(r.cells[0], label, size=7, bold=is_juara, color=col)
        set_cell(r.cells[1], val_f1, size=7, bold=is_juara, color=col, align_center=True)
        # Micro/W-F1/Acc tetap test (val tidak di-log)
        set_cell(r.cells[2], tmic, size=7, bold=is_juara, color=col, align_center=True)
        set_cell(r.cells[3], twfi, size=7, bold=is_juara, color=col, align_center=True)
        set_cell(r.cells[4], tacc, size=7, bold=is_juara, color=col, align_center=True)
        set_cell(r.cells[5], ww, size=7, bold=is_juara, color=col, align_center=True)

    # ── (4) Tighten all row heights so 28 rows fit ──
    # Existing table height ~3.80" → 28 rows × 0.135" ≈ 3.78"
    new_row_h = Emu(123000)  # ≈ 0.134"
    for row in tbl.rows:
        row.height = new_row_h

    # ── (5) Add footnote disclaimer below table ──
    footnote_top = 0.88 + (28 * 0.134) + 0.05  # right after table
    fn = sl.shapes.add_textbox(Inches(0.30), Inches(footnote_top),
                                Inches(4.90), Inches(0.30))
    fn.text_frame.word_wrap = True
    fn.text_frame.margin_left = Emu(0)
    fn.text_frame.margin_right = Emu(0)
    fn.text_frame.margin_top = Emu(0)
    fn.text_frame.margin_bottom = Emu(0)
    p = fn.text_frame.paragraphs[0]
    for run in list(p.runs):
        run._r.getparent().remove(run._r)
    p.text = ""
    run = p.add_run()
    run.text = ("* Micro / W-F1 / Acc / w = test set (val-only Micro/W-F1/Acc tidak di-log "
                "per config; w_best dipilih val-tuned).")
    run.font.size = Pt(6.5)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    print(f"   footnote added at top={footnote_top:.2f}\"")

    prs.save(PPTX)
    print(f"\nSaved: {PPTX}")


if __name__ == "__main__":
    main()
