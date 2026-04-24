"""
Update slide 133 SLIDE 36 Results table columns:
  Before: Config | Val F1 | Test F1 | Test Acc | w
  After:  Model  | Macro  | Micro   | W-F1     | Acc  | w
  (konsisten dengan Skema 1 tables, semua dari test set)

Val metric tetap disebut di notes tapi tidak di tabel — test metrics lebih
intuitif untuk presentation (konsisten dengan slide 122-125).
"""
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

PPTX   = Path('docs/PPT Bimbingan.pptx')
BACKUP = Path('docs/PPT Bimbingan_pre_cols.pptx')


def set_cell(cell, txt, size_pt=8, bold=False):
    cell.text = ''
    p = cell.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = txt
    run.font.size = Pt(size_pt)
    run.font.bold = bold


def main():
    if not PPTX.exists():
        raise FileNotFoundError(PPTX)
    shutil.copy2(PPTX, BACKUP)
    print(f'Backup: {BACKUP}')

    prs = Presentation(PPTX)
    slide = prs.slides[132]  # 1-idx 133 = 0-idx 132 (SLIDE 36 Results)

    # Find the 15-config table (should be first table, 16 rows × 5 cols existing)
    tables = [sh for sh in slide.shapes if sh.has_table]
    # Identify by size: first table at x=0.3 with 16 rows
    main_tbl = None
    for sh in tables:
        if len(sh.table.rows) == 16 and len(sh.table.columns) == 5:
            main_tbl = sh
            break
    if main_tbl is None:
        raise RuntimeError('15-config table not found in slide 133')

    # Remove the old table + re-insert new one with 6 columns
    old_left = main_tbl.left
    old_top = main_tbl.top
    old_width = main_tbl.width
    old_height = main_tbl.height
    main_tbl._element.getparent().remove(main_tbl._element)

    # New table: 16 rows × 6 cols
    new_tbl_shape = slide.shapes.add_table(
        rows=16, cols=6,
        left=old_left, top=old_top,
        width=old_width + Inches(0.4), height=old_height)
    t = new_tbl_shape.table

    # Column widths (tight)
    widths = [Inches(1.85), Inches(0.65), Inches(0.65), Inches(0.65), Inches(0.65), Inches(0.45)]
    for ci, w in enumerate(widths):
        t.columns[ci].width = w

    # Header + data (all test metrics, val tidak ditampilkan)
    header = ['Model', 'Macro', 'Micro', 'W-F1', 'Acc', 'w']
    rows_data = [
        ['FCNN B1',              '0.589', '0.741', '0.760', '0.741', '—'],
        ['FCNN B2',              '0.575', '0.702', '0.730', '0.702', '—'],
        ['FCNN B3',              '0.634', '0.749', '0.764', '0.749', '—'],
        ['CNN TL B1',            '0.634', '0.791', '0.800', '0.791', '—'],
        ['CNN TL B2',            '0.516', '0.581', '0.610', '0.581', '—'],
        ['CNN TL B3',            '0.705', '0.840', '0.850', '0.840', '—'],
        ['Intermediate TL B1',   '0.686', '0.790', '0.803', '0.790', '—'],
        ['Intermediate TL B2',   '0.665', '0.815', '0.819', '0.815', '—'],
        ['Intermediate TL B3',   '0.689', '0.828', '0.834', '0.828', '—'],
        ['Late Fusion TL B1',    '0.653', '0.797', '0.813', '0.797', '0.25'],
        ['Late Fusion TL B2',    '0.646', '0.774', '0.788', '0.774', '0.30'],
        ['Late Fusion TL B3 ⭐', '0.637', '0.784', '0.795', '0.784', '0.15'],
        ['Early Fusion TL B1',   '0.587', '0.783', '0.774', '0.783', '—'],
        ['Early Fusion TL B2',   '0.584', '0.650', '0.676', '0.650', '—'],
        ['Early Fusion TL B3',   '0.699', '0.820', '0.835', '0.820', '—'],
    ]

    for ci, h in enumerate(header):
        set_cell(t.cell(0, ci), h, size_pt=8.5, bold=True)
    for ri, row in enumerate(rows_data, start=1):
        for ci, val in enumerate(row):
            bold = '⭐' in row[0]
            set_cell(t.cell(ri, ci), val, size_pt=8, bold=bold)

    prs.save(PPTX)
    print(f'Saved: {PPTX}')


if __name__ == '__main__':
    main()
